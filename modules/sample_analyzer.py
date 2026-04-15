"""
샘플 정밀 분석기 — 구조 + 디자인 완전 수치화
각 샘플 PPT의 모든 도형, 표, 차트, 그룹의 디자인 속성을 추출하여 패턴 DB 저장
"""
import os
import json
import glob
import copy
from pathlib import Path
from pptx import Presentation
from pptx.util import Emu, Pt, Cm
from pptx.oxml.ns import qn
from lxml import etree


class SampleAnalyzer:
    def __init__(self, config):
        self.config = config
        self.base_dir = Path(__file__).parent.parent
        self.samples_dir = self.base_dir / config.get("samples_dir", "samples")
        self.patterns_file = self.base_dir / "sample_patterns.json"
        self.patterns = []
        
        if self.patterns_file.exists():
            with open(self.patterns_file, "r", encoding="utf-8") as f:
                self.patterns = json.load(f)
        else:
            self._analyze_all()
    
    def _analyze_all(self):
        self.patterns = []
        for f in glob.glob(str(self.samples_dir / "*.pptx")):
            pattern = self.analyze_single(f)
            if pattern:
                self.patterns.append(pattern)
        self.save_patterns()
    
    def analyze_single(self, file_path):
        """단일 PPT 정밀 분석 — 구조 + 디자인 + 데이터 특성"""
        try:
            prs = Presentation(file_path)
        except:
            return None
        
        basename = os.path.basename(file_path)
        slides_info = []
        
        # 전체 특성
        features = {
            "categories": 0, "instructors": 0, "has_modules": False,
            "tables": 0, "charts": 0, "groups": 0,
        }
        counts = {"exec_slides": 0, "quant_slides": 0, "qual_slides": 0}
        
        for idx, slide in enumerate(prs.slides):
            slide_data = self._analyze_slide(slide, idx)
            slides_info.append(slide_data)
            
            role = slide_data["role"]
            if role == "exec_summary": counts["exec_slides"] += 1
            if role == "quant": counts["quant_slides"] += 1
            if role == "qual": counts["qual_slides"] += 1
            features["tables"] += len(slide_data["tables"])
            features["charts"] += len(slide_data["charts"])
            features["groups"] += len(slide_data["groups"])
            
            # Executive Summary 표에서 카테고리/강사 수 추출
            for tbl in slide_data["tables"]:
                if role == "exec_summary" and tbl["rows"] == 3:
                    features["categories"] = max(features["categories"], tbl["cols"])
                    for header in tbl.get("headers", []):
                        if "강사" in header:
                            features["instructors"] += 1
                        if "모듈" in header.lower() or "[L" in header:
                            features["has_modules"] = True
        
        structure = [s["role"] for s in slides_info]
        
        return {
            "name": basename,
            "path": str(file_path),
            "slide_count": len(prs.slides),
            "features": features,
            "counts": counts,
            "structure": structure,
            "slides": slides_info,
        }
    
    def _analyze_slide(self, slide, idx):
        """슬라이드 정밀 분석"""
        layout_name = slide.slide_layout.name
        
        # 역할 판별
        title = ""
        texts = []
        tables = []
        charts = []
        groups = []
        shapes_info = []
        
        for shape in slide.shapes:
            shape_data = {
                "name": shape.name,
                "type": str(shape.shape_type),
                "left": shape.left,
                "top": shape.top,
                "width": shape.width,
                "height": shape.height,
            }
            
            if shape.has_text_frame:
                shape_data["text"] = shape.text_frame.text[:100]
                for p in shape.text_frame.paragraphs:
                    t = p.text.strip()
                    if t:
                        texts.append(t)
                        if not title and 2 < len(t) < 40:
                            title = t
                # 텍스트 스타일
                if shape.text_frame.paragraphs and shape.text_frame.paragraphs[0].runs:
                    run = shape.text_frame.paragraphs[0].runs[0]
                    shape_data["font"] = self._safe_font_info(run)
            
            if shape.has_table:
                tbl = shape.table
                tbl_data = self._analyze_table(tbl, shape)
                tables.append(tbl_data)
                shape_data["table"] = tbl_data
            
            if hasattr(shape, 'has_chart') and shape.has_chart:
                chart_data = self._analyze_chart(shape.chart, shape)
                charts.append(chart_data)
                shape_data["chart"] = chart_data
            
            if shape.shape_type == 6:  # Group
                grp_data = self._analyze_group(shape)
                groups.append(grp_data)
                shape_data["group"] = grp_data
            
            shapes_info.append(shape_data)
        
        role = self._determine_role(title, layout_name, tables, charts, groups)
        
        return {
            "index": idx,
            "layout": layout_name,
            "role": role,
            "title": title,
            "shapes": shapes_info,
            "tables": tables,
            "charts": charts,
            "groups": groups,
        }
    
    def _analyze_table(self, table, shape):
        """표 정밀 분석 — 셀 스타일, 색상, 폰트, 크기"""
        rows = len(table.rows)
        cols = len(table.columns)
        
        col_widths = [table.columns[c].width for c in range(cols)]
        row_heights = [table.rows[r].height for r in range(rows)]
        
        # 헤더 행 스타일
        headers = []
        header_style = {}
        if rows > 0:
            for c in range(cols):
                cell = table.cell(0, c)
                headers.append(cell.text[:30])
                if c == 0:
                    header_style = self._extract_cell_style(cell)
        
        # 데이터 행 스타일 (행 1)
        data_style = {}
        if rows > 1:
            data_style = self._extract_cell_style(table.cell(1, 0))
        
        # 표 전체 스타일 ID
        tblPr = table._tbl.find(qn('a:tblPr'))
        style_id = ""
        if tblPr is not None:
            style_el = tblPr.find(qn('a:tblStyle'))
            if style_el is not None:
                style_id = style_el.get('val', '')
        
        # 데이터 내용 (미리보기용)
        data_rows = []
        for r in range(min(rows, 5)):
            row_data = [table.cell(r, c).text[:30] for c in range(cols)]
            data_rows.append(row_data)
        
        return {
            "rows": rows, "cols": cols,
            "col_widths": col_widths,
            "row_heights": row_heights,
            "headers": headers,
            "header_style": header_style,
            "data_style": data_style,
            "style_id": style_id,
            "position": {"left": shape.left, "top": shape.top,
                         "width": shape.width, "height": shape.height},
            "data_rows": data_rows,
        }
    
    def _extract_cell_style(self, cell):
        """셀 스타일 추출"""
        style = {}
        tc = cell._tc
        tcPr = tc.find(qn('a:tcPr'))
        if tcPr is not None:
            sf = tcPr.find(qn('a:solidFill'))
            if sf is not None:
                clr = sf.find(qn('a:srgbClr'))
                if clr is not None:
                    style["bg_color"] = clr.get('val')
        
        # 폰트
        if cell.text_frame.paragraphs and cell.text_frame.paragraphs[0].runs:
            run = cell.text_frame.paragraphs[0].runs[0]
            style.update(self._safe_font_info(run))
        
        # 정렬
        if cell.text_frame.paragraphs:
            algn = cell.text_frame.paragraphs[0].alignment
            style["alignment"] = str(algn) if algn else None
        
        return style
    
    def _analyze_chart(self, chart, shape):
        """차트 스타일 분석"""
        chart_data = {
            "type": str(chart.chart_type),
            "position": {"left": shape.left, "top": shape.top,
                         "width": shape.width, "height": shape.height},
        }
        
        # 시리즈 색상
        series_colors = []
        try:
            for ser in chart.series:
                ser_el = ser._element
                spPr = ser_el.find(qn('c:spPr'))
                if spPr is not None:
                    sf = spPr.find(qn('a:solidFill'))
                    if sf is not None:
                        clr = sf.find(qn('a:srgbClr'))
                        if clr is not None:
                            series_colors.append(clr.get('val'))
        except:
            pass
        chart_data["series_colors"] = series_colors
        
        # 데이터 포인트 수
        try:
            chart_data["data_points"] = len(list(chart.series[0].values))
        except:
            chart_data["data_points"] = 0
        
        return chart_data
    
    def _analyze_group(self, shape):
        """그룹 도형 분석"""
        grp_data = {
            "position": {"left": shape.left, "top": shape.top,
                         "width": shape.width, "height": shape.height},
            "children": [],
        }
        try:
            for child in shape.shapes:
                child_info = {
                    "name": child.name,
                    "type": str(child.shape_type),
                    "left": child.left, "top": child.top,
                    "width": child.width, "height": child.height,
                }
                if child.has_text_frame:
                    child_info["text"] = child.text_frame.text[:60]
                    if child.text_frame.paragraphs and child.text_frame.paragraphs[0].runs:
                        run = child.text_frame.paragraphs[0].runs[0]
                        child_info["font"] = {
                            "name": run.font.name,
                            "size": int(run.font.size) if run.font.size else None,
                            "bold": run.font.bold,
                        }
                # 채우기 색상
                spPr = child._element.find(qn('p:spPr'))
                if spPr is not None:
                    sf = spPr.find(qn('a:solidFill'))
                    if sf is not None:
                        clr = sf.find(qn('a:srgbClr'))
                        if clr is not None:
                            child_info["fill_color"] = clr.get('val')
                grp_data["children"].append(child_info)
        except:
            pass
        return grp_data
    
    def _determine_role(self, title, layout, tables, charts, groups):
        """슬라이드 역할 판별"""
        if "앞표지" in layout: return "cover"
        if "목차" in layout: return "toc"
        if "뒤표지" in layout: return "back"
        if "Executive" in title: return "exec_summary"
        if "정량" in title: return "quant"
        if "정성" in title: return "qual"
        if "과정 개요" in title: return "overview"
        if "일정" in title: return "schedule"
        if "만족도 결과" in title: return "section_satisfaction"
        if "Appendix" in title: return "appendix"
        if "실시 개요" in title or "실시개요" in title: return "section_overview"
        if not title and "빈 화면" in layout: return "photo"
        if not title and "캡션" in layout: return "section_divider"
        return "other"
    
    def find_best_match(self, data_summary):
        """데이터 특성에 가장 잘 맞는 샘플 찾기"""
        if not self.patterns:
            self._analyze_all()
        
        best = None
        best_score = -1
        company = data_summary.get("company", "").strip()
        
        # 같은 고객사 샘플이 있으면 최우선 선택
        if company:
            for pattern in self.patterns:
                if company in pattern.get("name", ""):
                    return pattern
        
        for pattern in self.patterns:
            score = 0
            feat = pattern.get("features", {})
            counts = pattern.get("counts", {})
            
            # 카테고리 수 유사도 (최우선 — 카테고리 수가 보고서 구조를 결정)
            target_cats = data_summary.get("categories", 0)
            cat_diff = abs(feat.get("categories", 0) - target_cats)
            score += max(0, 20 - cat_diff * 5)
            
            # 강사 수 유사도 (10명 이상은 스캔 오류)
            sample_instr = min(feat.get("instructors", 0), 5)
            target_instr = data_summary.get("num_instructors", 0)
            instr_diff = abs(sample_instr - target_instr)
            score += max(0, 6 - instr_diff * 2)
            
            # 모듈 유무 일치
            if feat.get("has_modules") == data_summary.get("has_modules"):
                score += 8
            
            # 구조 복잡도 매칭
            if counts.get("exec_slides", 0) > 1:
                score += 4
            if counts.get("quant_slides", 0) > 1 and data_summary.get("has_modules"):
                score += 4
            if counts.get("qual_slides", 0) >= 1:
                score += 2
            
            # 전체 피처 거리 (가까울수록 보너스)
            total_diff = cat_diff + instr_diff
            score += max(0, 5 - total_diff)
            
            if score > best_score:
                best_score = score
                best = pattern
        
        return best
    
    def get_patterns(self):
        if not self.patterns:
            self._analyze_all()
        return self.patterns
    
    def save_patterns(self):
        # JSON 직렬화 가능하도록 Emu 값들을 int로 변환
        clean = json.loads(json.dumps(self.patterns, default=str))
        with open(self.patterns_file, "w", encoding="utf-8") as f:
            json.dump(clean, f, ensure_ascii=False, indent=2)
    @staticmethod
    def _safe_font_info(run):
        """폰트 정보 안전 추출"""
        info = {
            "name": run.font.name,
            "size": int(run.font.size) if run.font.size else None,
            "bold": run.font.bold,
            "color": None,
        }
        try:
            if run.font.color and run.font.color.type is not None:
                info["color"] = str(run.font.color.rgb)
        except:
            pass
        return info
