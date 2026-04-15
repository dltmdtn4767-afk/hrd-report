"""
PPT 미리보기 엔진 — PPT를 HTML 슬라이드 데이터로 변환
LibreOffice 불필요, python-pptx 직접 파싱
"""
from pptx import Presentation
from pptx.util import Emu, Pt


def generate_preview(pptx_path):
    """PPT 파일 → 슬라이드별 미리보기 데이터 리스트"""
    prs = Presentation(pptx_path)
    slides = []
    
    for idx, slide in enumerate(prs.slides):
        layout_name = slide.slide_layout.name
        texts = []
        tables = []
        charts = []
        groups = []
        title = ""
        table_rows = 0
        group_count = 0
        
        for shape in slide.shapes:
            if shape.has_text_frame:
                for para in shape.text_frame.paragraphs:
                    t = para.text.strip()
                    if t:
                        texts.append(t)
                        if not title and 2 < len(t) < 40:
                            title = t
            
            if shape.has_table:
                tbl = shape.table
                rows = []
                for r in range(len(tbl.rows)):
                    cells = [tbl.cell(r, c).text[:30] for c in range(len(tbl.columns))]
                    rows.append(cells)
                tables.append({
                    "name": shape.name,
                    "size": f"{len(tbl.rows)}x{len(tbl.columns)}",
                    "rows": rows[:20],  # 최대 20행
                })
                table_rows += len(tbl.rows) - 1  # 헤더 제외
            
            if hasattr(shape, 'has_chart') and shape.has_chart:
                chart = shape.chart
                try:
                    series_data = list(chart.series[0].values)[:20] if chart.series else []
                    charts.append({
                        "type": str(chart.chart_type),
                        "values": [round(v, 2) for v in series_data],
                    })
                except:
                    charts.append({"type": "chart", "values": []})
            
            if shape.shape_type == 6:
                grp_texts = []
                try:
                    for child in shape.shapes:
                        if child.has_text_frame and child.text_frame.text.strip():
                            grp_texts.append(child.text_frame.text[:60])
                except:
                    pass
                if grp_texts:
                    groups.append(grp_texts)
                    group_count += 1
        
        slides.append({
            "index": idx + 1,
            "layout": layout_name,
            "title": title,
            "texts": texts[:10],
            "tables": tables,
            "charts": charts,
            "groups": groups,
            "table_rows": table_rows,
            "group_count": group_count,
        })
    
    return slides
