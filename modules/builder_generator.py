"""
builder_generator.py
빌더 구성 → 실제 템플릿 기반 PPT 생성
"""
from pathlib import Path
from datetime import datetime
import copy

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.chart.data import ChartData
from pptx.enum.chart import XL_CHART_TYPE
from pptx.enum.text import MSO_ANCHOR
from modules.ppt_constants import (
    FONT_NAME, BRAND_COLORS,
    CHART_SIZE, TABLE_SIZE,
    CHART_CONFIG
)

BASE_DIR = Path(__file__).parent.parent
TEMPLATE_PATH = BASE_DIR / "templates" / "template.pptx"
OUTPUT_DIR = BASE_DIR / "output"
OUTPUT_DIR.mkdir(exist_ok=True)

# 템플릿 슬라이드 인덱스 (0-based)
TPL_COVER    = 0
TPL_TOC      = 1
TPL_SEC1     = 2   # Ⅰ 과정 개요
TPL_OVERVIEW = 3   # 과정 개요 표
TPL_SCHEDULE = 4   # 교육 일정표
TPL_SEC2     = 7   # Ⅱ 만족도 결과
TPL_SUMMARY  = 8   # Executive Summary
TPL_QUANT    = 9   # 정량 평가 결과 (복사용)
TPL_QUAL     = 10  # 정성 평가 결과
TPL_BACK     = 11  # 뒤표지


# ────────────────────────────────────────────
# 유틸
# ────────────────────────────────────────────
def _set_text(shape, text: str, bold=None, size=None, color=None, align=None):
    """텍스트박스/플레이스홀더 텍스트 교체"""
    if not shape.has_text_frame:
        return
    tf = shape.text_frame
    # 전체 텍스트 첫 단락에 설정
    for i, para in enumerate(tf.paragraphs):
        for run in para.runs:
            run.font.name = FONT_NAME
            if bold is not None: run.font.bold = bold
            if size: run.font.size = Pt(size)
            if color: run.font.color.rgb = color # color is already RGBColor or from constants
            break
        if i == 0 and not para.runs:
            run = para.add_run()
            run.text = text
            run.font.name = FONT_NAME
            if bold is not None: run.font.bold = bold
            if size: run.font.size = Pt(size)
            if color: run.font.color.rgb = color
        if align and i == 0:
            para.alignment = align
        if i > 0:
            for run in para.runs:
                run.text = ''


def _find_shape(slide, name_contains: str):
    for sh in slide.shapes:
        if name_contains.lower() in sh.name.lower():
            return sh
    return None


def _set_table_cell(table, row, col, text, bold=False, bg=None):
    try:
        cell = table.cell(row, col)
        tf = cell.text_frame
        for p in tf.paragraphs:
            for r in p.runs:
                r.text = ''
            run = tf.paragraphs[0].add_run()
            run.text = str(text)
            run.font.name = FONT_NAME
            run.font.bold = bold
            from pptx.enum.text import MSO_VERTICAL_ANCHOR
            cell.vertical_anchor = MSO_VERTICAL_ANCHOR.MIDDLE # 수직 중앙 정렬
            if bg:
                from pptx.oxml.ns import qn
                from lxml import etree
                # fill solid
                tc = cell._tc
                tcPr = tc.find(qn('a:tcPr')) or etree.SubElement(tc, qn('a:tcPr'))
                solidFill = etree.SubElement(tcPr, qn('a:solidFill'))
                srgb = etree.SubElement(solidFill, qn('a:srgbClr'))
                srgb.set('val', bg.lstrip('#'))
    except Exception:
        pass


def _clone_slide(prs: Presentation, src_idx: int) -> object:
    """슬라이드 복제 (XML 복사)"""
    import copy
    from pptx.oxml.ns import qn
    from lxml import etree

    src_slide = prs.slides[src_idx]
    blank_layout = src_slide.slide_layout

    new_slide = prs.slides.add_slide(blank_layout)

    # XML 내용 복사
    new_slide.shapes._spTree.clear()
    for elem in copy.deepcopy(src_slide.shapes._spTree):
        new_slide.shapes._spTree.append(elem)

    return new_slide


def merge_table_headers(table):
    """1열(인덱스 0)의 동일한 텍스트를 가진 셀을 찾아 자동 병합"""
    if len(table.rows) <= 1:
        return
        
    start_row = 1 # 0행은 타이틀 헤더이므로 제외
    current_text = table.cell(start_row, 0).text

    for i in range(2, len(table.rows)):
        cell_text = table.cell(i, 0).text
        if cell_text != current_text:
            if i - 1 > start_row:
                table.cell(start_row, 0).merge(table.cell(i - 1, 0))
            start_row = i
            current_text = cell_text
            
    # 마지막 그룹 병합 처리 (끝까지 같은 내용일 경우)
    if len(table.rows) - 1 > start_row:
        table.cell(start_row, 0).merge(table.cell(len(table.rows) - 1, 0))


def apply_table_style(table):
    """테이블 전체 '나눔바른고딕 Light' 폰트 및 세로 중앙 정렬"""
    for row in table.rows:
        for cell in row.cells:
            cell.vertical_anchor = MSO_ANCHOR.MIDDLE # 세로 중앙 정렬
            for paragraph in cell.text_frame.paragraphs:
                for run in paragraph.runs:
                    run.font.name = FONT_NAME


# ────────────────────────────────────────────
# 슬라이드 타입별 채우기 함수
# ────────────────────────────────────────────
def _fill_cover(slide, data: dict):
    """앞표지: 고객사명 + 과정명"""
    for sh in slide.shapes:
        if sh.has_text_frame:
            txt = sh.text_frame.text
            if '고객사명' in txt:
                _set_text(sh, data.get('company', ''))
            elif '과정명' in txt:
                _set_text(sh, data.get('course', ''))


def _fill_overview(slide, data: dict):
    """과정 개요 표"""
    rows = data.get('rows', [])
    tbl_shape = _find_shape(slide, '표')
    if not tbl_shape or not tbl_shape.has_table:
        return
    tbl = tbl_shape.table
    for i, row in enumerate(rows):
        try:
            if i < tbl.rows.__len__():
                _set_table_cell(tbl, i, 0, row.get('key', ''), bold=True)
                _set_table_cell(tbl, i, 1, row.get('val', ''))
        except Exception:
            pass


def _fill_schedule(slide, data: dict):
    """교육 일정표"""
    rows = data.get('rows', [])
    tbl_shape = _find_shape(slide, '표')
    if not tbl_shape or not tbl_shape.has_table:
        return
    tbl = tbl_shape.table
    for i, row in enumerate(rows):
        try:
            r_idx = i + 1  # 헤더 제외
            if r_idx < tbl.rows.__len__():
                _set_table_cell(tbl, r_idx, 0, row.get('label', ''))
                _set_table_cell(tbl, r_idx, 1, row.get('date', ''))
                _set_table_cell(tbl, r_idx, 2, row.get('place', ''))
                _set_table_cell(tbl, r_idx, 3, row.get('count', ''))
        except Exception:
            pass


def _fill_summary(slide, data: dict):
    """Executive Summary — 차트 + 스타일A 표 (가로 영역별)"""
    cats = data.get('categories', [])
    overall = data.get('overall', 0)

    # 차트 업데이트
    for sh in slide.shapes:
        try:
            chart = sh.chart
            cd = ChartData()
            cd.categories = [c.get('name', '') for c in cats]
            cd.add_series('평균', [round(c.get('avg', 0), 2) for c in cats])
            chart.replace_data(cd)
            break
        except Exception:
            pass

    # 표 — 스타일A: 카테고리별 열 + 과정 전반 하이라이트
    tbl_shape = _find_shape(slide, '표')
    if tbl_shape and tbl_shape.has_table:
        tbl = tbl_shape.table
        n_cols = len(tbl.columns)
        try:
            for ci, cat in enumerate(cats):
                if ci < n_cols:
                    _set_table_cell(tbl, 0, ci, cat.get('name', ''), bold=True)
                    _set_table_cell(tbl, 1, ci, f"{cat.get('avg', 0):.2f}")
            # 마지막 열 = 과정 전반 (하이라이트)
            last_col = min(len(cats), n_cols - 1)
            if last_col < n_cols:
                _set_table_cell(tbl, 0, last_col, '과정 전반', bold=True, bg='F0BEF8')
                _set_table_cell(tbl, 1, last_col, f"{overall:.2f}")
            # 평균 행
            if tbl.rows.__len__() > 2:
                _set_table_cell(tbl, 2, 0, f"{overall:.2f}", bold=True, bg='E7E6E6')
        except Exception:
            pass


def _fill_quant_chart(slide, group: dict):
    """정량 차트 슬라이드 — 네이티브 차트(데이터시트 포함) + 네이티브 표 생성
    PPT에서 차트 더블클릭 → 데이터 편집 가능, 표/차트 모두 복사 붙여넣기 가능
    """
    questions = group.get('questions', [])
    title     = group.get('title', '정량 평가 결과')

    # 슬라이드 제목
    title_sh = _find_shape(slide, '제목') or _find_shape(slide, 'TextBox 224')
    if title_sh:
        _set_text(title_sh, title, bold=True)

    # ── 기존 표가 있으면 템플릿 스타일B로 채우기 ──
    tbl_shape = _find_shape(slide, '표')
    if tbl_shape and tbl_shape.has_table:
        tbl = tbl_shape.table
        # 헤더
        _set_table_cell(tbl, 0, 0, '항목', bold=True, bg='F1F5F9')
        _set_table_cell(tbl, 0, 1, '문항', bold=True, bg='F1F5F9')
        _set_table_cell(tbl, 0, 2, '평균', bold=True, bg='F1F5F9')
        if len(tbl.columns) > 3:
            _set_table_cell(tbl, 0, 3, '응답인원', bold=True, bg='F1F5F9')

        for i, q in enumerate(questions):
            r = i + 1
            if r >= tbl.rows.__len__(): break
            
            _set_table_cell(tbl, r, 0, q.get('category', q.get('id', '')))
            _set_table_cell(tbl, r, 1, q.get('label', ''))
            avg_val = q.get('avg', 0)
            _set_table_cell(tbl, r, 2, f"{avg_val:.02f}")
            if len(tbl.columns) > 3:
                _set_table_cell(tbl, r, 3, str(q.get('count', '')))
        
        # 스타일 및 병합 적용
        apply_table_style(tbl)
        merge_table_headers(tbl)

    # ── 네이티브 차트 삽입 (데이터 시트 포함 → PPT에서 더블클릭 편집 가능) ──
    has_chart = False
    for sh in slide.shapes:
        try:
            chart = sh.chart
            has_chart = True
            cd = ChartData()
            cd.categories = [q.get('label', str(i+1))[:20] for i, q in enumerate(questions)]
            cd.add_series('평균', [round(q.get('avg', 0), 2) for q in questions])
            chart.replace_data(cd)
            # 차트 스타일 설정
            try:
                plot = chart.plots[0]
                plot.gap_width = CHART_CONFIG["gap_width"]
                chart.category_axis.tick_labels.font.name = FONT_NAME
                chart.value_axis.tick_labels.font.name = FONT_NAME
                chart.value_axis.major_gridlines.format.line.color.rgb = BRAND_COLORS["GRID"]
                for pt_idx, q in enumerate(questions):
                    try:
                        pt = plot.series[0].points[pt_idx]
                        avg = q.get('avg', 0)
                        if avg >= 4.5:
                            pt.format.fill.solid()
                            pt.format.fill.fore_color.rgb = BRAND_COLORS["SUCCESS"]
                        elif avg < 3.5:
                            pt.format.fill.solid()
                            pt.format.fill.fore_color.rgb = BRAND_COLORS["DANGER"]
                        else:
                            pt.format.fill.solid()
                            pt.format.fill.fore_color.rgb = BRAND_COLORS["BLUE"]
                    except Exception:
                        pass
            except Exception:
                pass
            break
        except Exception:
            pass

    # 차트가 없으면 새로 추가
    if not has_chart and questions:
        _add_native_chart(slide, questions, title)


def _add_native_chart(slide, questions, title):
    """PPT 네이티브 막대 차트를 슬라이드에 새로 추가 (데이터 시트 포함)"""
    cd = ChartData()
    labels = [q.get('label', '')[:20] for q in questions]
    values = [round(q.get('avg', 0), 2) for q in questions]
    cd.categories = labels
    cd.add_series('평균', values)

    # 차트 위치: 슬라이드 왼쪽 상단 ~ 중앙
    x = Inches(0.5)
    y = Inches(1.2)
    cx = CHART_SIZE["width"]
    cy = CHART_SIZE["height"]

    chart_frame = slide.shapes.add_chart(
        XL_CHART_TYPE.COLUMN_CLUSTERED,
        x, y, cx, cy, cd
    )
    chart = chart_frame.chart

    # 차트 스타일
    chart.has_legend = False
    try:
        plot = chart.plots[0]
        plot.gap_width = CHART_CONFIG["gap_width"]
        # 데이터 레이블 표시
        plot.has_data_labels = True
        data_labels = plot.data_labels
        data_labels.number_format = CHART_CONFIG["label_format"]
        data_labels.font.size = Pt(9)
        data_labels.font.bold = True
        data_labels.font.name = FONT_NAME

        # 색상: 점수에 따라 분류
        for pt_idx, q in enumerate(questions):
            try:
                pt = plot.series[0].points[pt_idx]
                avg = q.get('avg', 0)
                pt.format.fill.solid()
                if avg >= 4.5:
                    pt.format.fill.fore_color.rgb = BRAND_COLORS["SUCCESS"]
                elif avg < 3.5:
                    pt.format.fill.fore_color.rgb = BRAND_COLORS["DANGER"]
                else:
                    pt.format.fill.fore_color.rgb = BRAND_COLORS["BLUE"]
            except Exception:
                pass
    except Exception:
        pass

    # Y축 범위
    try:
        value_axis = chart.value_axis
        value_axis.minimum_scale = 0
        value_axis.maximum_scale = 5
        value_axis.major_unit = 1
        value_axis.has_title = False
    except Exception:
        pass

    # X축 폰트
    try:
        cat_axis = chart.category_axis
        cat_axis.tick_labels.font.size = Pt(8)
        cat_axis.tick_labels.font.name = FONT_NAME
        chart.value_axis.tick_labels.font.name = FONT_NAME
        chart.value_axis.major_gridlines.format.line.color.rgb = BRAND_COLORS["GRID"]
    except Exception:
        pass



def _fill_qual(slide, qual_data: list):
    """정성 평가 슬라이드 — 텍스트 그룹 채우기"""
    if not qual_data:
        return
    # 각 Q별 공통응답을 텍스트로 합침
    lines = []
    for oe in qual_data:
        qid   = oe.get('id', '')
        label = oe.get('label', '')
        lines.append(f"■ {qid} {label}")
        for g in oe.get('groups', []):
            lines.append(f"  · {g.get('label','')} ({g.get('count',0)}건)")
        lines.append('')

    full_text = '\n'.join(lines)

    # 텍스트박스에 삽입 (첫 번째 큰 텍스트박스)
    for sh in slide.shapes:
        if sh.has_text_frame and sh.width > Inches(4):
            _set_text(sh, full_text)
            break


def _fill_custom_quant(slide, custom_data: dict):
    """커스텀 그룹 슬라이드 — 그룹별 평균 차트 + 표"""
    groups = custom_data.get('groups', [])
    title = custom_data.get('title', '커스텀 분석')
    table_style = custom_data.get('tableStyle', 'B')

    # 슬라이드 제목
    title_sh = _find_shape(slide, '제목') or _find_shape(slide, 'TextBox')
    if title_sh:
        _set_text(title_sh, title, bold=True)

    # 차트: 있으면 데이터 교체, 없으면 네이티브 차트 추가
    has_chart = False
    for sh in slide.shapes:
        try:
            chart = sh.chart
            has_chart = True
            cd = ChartData()
            cd.categories = [g.get('name', '') for g in groups]
            cd.add_series('평균', [round(g.get('avg', 0), 2) for g in groups])
            chart.replace_data(cd)
            break
        except Exception:
            pass

    if not has_chart and groups:
        # 그룹을 문항처럼 변환해서 네이티브 차트 추가
        fake_qs = [{"label": g.get("name",""), "avg": g.get("avg",0)} for g in groups]
        _add_native_chart(slide, fake_qs, title)

    # 표 채우기
    tbl_shape = _find_shape(slide, '표')
    if not tbl_shape or not tbl_shape.has_table:
        return
    tbl = tbl_shape.table

    if table_style == 'A':
        n_cols = len(tbl.columns)
        for ci, g in enumerate(groups):
            if ci < n_cols:
                _set_table_cell(tbl, 0, ci, g.get('name', ''), bold=True)
                _set_table_cell(tbl, 1, ci, f"{g.get('avg', 0):.2f}")
        overall = sum(g.get('avg', 0) for g in groups) / max(len(groups), 1)
        last = min(len(groups), n_cols - 1)
        if last < n_cols:
            _set_table_cell(tbl, 0, last, '과정 전반', bold=True, bg='EFF6FF')
            _set_table_cell(tbl, 1, last, f"{overall:.2f}")
    else:
        _set_table_cell(tbl, 0, 0, '항목', bold=True, bg='F1F5F9')
        _set_table_cell(tbl, 0, 1, '문항', bold=True, bg='F1F5F9')
        _set_table_cell(tbl, 0, 2, '평균', bold=True, bg='F1F5F9')
        _set_table_cell(tbl, 0, 3, '응답인원', bold=True, bg='F1F5F9')
        row_idx = 1
        for g in groups:
            cur_cat_name = g.get('name', '')
            if row_idx < tbl.rows.__len__():
                _set_table_cell(tbl, row_idx, 0, cur_cat_name, bold=True, bg='F1F5F9')
                _set_table_cell(tbl, row_idx, 2, f"{g.get('avg', 0):.2f}")
                row_idx += 1
            for q in g.get('questions', []):
                if row_idx < tbl.rows.__len__():
                    _set_table_cell(tbl, row_idx, 0, cur_cat_name)
                    _set_table_cell(tbl, row_idx, 1, q.get('label', ''))
                    _set_table_cell(tbl, row_idx, 2, f"{q.get('avg', 0):.2f}")
                    _set_table_cell(tbl, row_idx, 3, str(q.get('count', '')))
                    row_idx += 1
        
        apply_table_style(tbl)
        merge_table_headers(tbl)


# ────────────────────────────────────────────
# 메인 빌드 함수
# ────────────────────────────────────────────
def build_custom_ppt(slides: list, quant_groups: list, qual_data: list,
                     data: dict, config: dict, template_path: str = None) -> str:
    tpl = Path(template_path) if template_path and Path(template_path).exists() else TEMPLATE_PATH
    if not tpl.exists():
        raise FileNotFoundError(f"템플릿 없음: {tpl}")
    prs = Presentation(str(tpl))

    # ── 1. 템플릿 슬라이드 조작 ──────────────────
    # 슬라이드 순서를 빌더 구성대로 재구성
    # 전략: 템플릿을 기반으로 각 슬라이드 타입의 내용만 채움

    # 표지 채우기
    cover_data = next((s['data'] for s in slides if s['type'] == 'cover'), {})
    _fill_cover(prs.slides[TPL_COVER], cover_data)

    # 과정 개요 표
    ov_data = next((s['data'] for s in slides if s['type'] == 'overview'), {})
    _fill_overview(prs.slides[TPL_OVERVIEW], ov_data)

    # 교육 일정표
    sch_data = next((s['data'] for s in slides if s['type'] == 'schedule'), {})
    _fill_schedule(prs.slides[TPL_SCHEDULE], sch_data)

    # Executive Summary
    summary_data = next((s['data'] for s in slides if s['type'] == 'summary'), {})
    if not summary_data.get('categories'):
        summary_data['categories'] = data.get('categories', [])
    _fill_summary(prs.slides[TPL_SUMMARY], summary_data)

    # 정성 평가
    _fill_qual(prs.slides[TPL_QUAL], qual_data)

    # ── 2. 정량 그룹 슬라이드 (슬라이드 10 복제) ──
    if quant_groups:
        _fill_quant_chart(prs.slides[TPL_QUANT], quant_groups[0])
        for g in quant_groups[1:]:
            new_slide = _clone_slide(prs, TPL_QUANT)
            _fill_quant_chart(new_slide, g)
            xml_slides = prs.slides._sldIdLst
            last = xml_slides[-1]
            qual_xml_idx = len(prs.slides) - 2
            xml_slides.remove(last)
            xml_slides.insert(qual_xml_idx, last)

    # ── 2b. 커스텀 그룹 슬라이드 ──
    custom_quant_slides = [s for s in slides if s.get('type') == 'custom_quant']
    for cs in custom_quant_slides:
        cs_data = cs.get('data', {})
        new_slide = _clone_slide(prs, TPL_QUANT)
        _fill_custom_quant(new_slide, cs_data)
        # 정성 슬라이드 앞에 삽입
        xml_slides = prs.slides._sldIdLst
        last = xml_slides[-1]
        qual_xml_idx = len(prs.slides) - 2
        xml_slides.remove(last)
        xml_slides.insert(qual_xml_idx, last)

    # ── 3. 파일명 & 저장 ─────────────────────────
    company  = cover_data.get('company', '')
    course   = cover_data.get('course', '')
    today    = datetime.now().strftime('%Y%m%d')
    filename = f"[결과보고서] {company}_{course}_{today}.pptx"
    out_path = OUTPUT_DIR / filename
    prs.save(str(out_path))
    return str(out_path)
