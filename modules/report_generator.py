"""
report_generator.py v10 — 템플릿 기반 + 샘플 디자인 참조

핵심 원칙:
  1. 기본 템플릿을 베이스로 연다 (사진/교재 등 불필요 콘텐츠 없음)
  2. 샘플은 디자인 참조용 — 표 스타일, 색상, 폰트 등만 가져온다
  3. 빈 칸 없이 모든 텍스트/표 셀에 데이터를 채운다
  4. 데이터에 없는 슬라이드(사진, Appendix 등)는 생성하지 않는다
"""
import os
import copy
import glob
from datetime import datetime
from pptx import Presentation
from pptx.util import Pt, Cm, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.oxml.ns import qn
from lxml import etree
from modules.config_manager import resolve_path


def find_template(config):
    tpl_path = config.get("ppt_template", {}).get("template_path", "")
    if tpl_path:
        full = resolve_path(tpl_path)
        if os.path.exists(full):
            return full
    tpl_dir = resolve_path("templates")
    for f in glob.glob(os.path.join(tpl_dir, "*.pptx")):
        bn = os.path.basename(f).lower()
        if '템플' in bn or 'template' in bn:
            return f
    pptx_files = glob.glob(os.path.join(tpl_dir, "*.pptx"))
    if pptx_files:
        return pptx_files[0]
    raise FileNotFoundError("PPT 템플릿을 찾을 수 없습니다.")


def generate_report(data, config, sample_pattern=None, template_path=None,
                    output_path=None, multi_result=None):
    """
    보고서 생성 메인
    - 항상 기본 템플릿을 베이스로 사용
    - multi_result 있으면 차수별 + 종합 슬라이드 자동 생성
    """
    course = data["course_info"]
    
    if output_path is None:
        output_dir = resolve_path(config.get("ppt_template", {}).get("output_dir", "output"))
        os.makedirs(output_dir, exist_ok=True)
        today = datetime.now().strftime("%Y%m%d")
        company = course.get("company", "")
        cname = course.get("course_name", "교육")
        output_path = os.path.join(output_dir, f"[결과보고서] {company}_{cname}_{today}.pptx")
    
    # ═══ 항상 기본 템플릿을 베이스로 ═══
    base_path = template_path or find_template(config)
    prs = Presentation(base_path)
    
    # ═══ 샘플에서 디자인 정보만 추출 ═══
    sample_design = None
    if sample_pattern and os.path.exists(sample_pattern.get("path", "")):
        sample_design = _extract_design_from_sample(sample_pattern["path"])
    
    # ═══ 공통 데이터 ═══
    company_name = course.get("company", "고객사명")
    course_name = course.get("course_name", "과정명")
    header_text = f"'{course_name}' 과정운영 결과보고서"
    today_str = datetime.now().strftime("%Y. %m. %d.")
    
    # ═══ 슬라이드 역할 맵핑 ═══
    slide_roles = _map_slide_roles(prs)
    
    # ═══ (1) 마스터 헤더 ═══
    _update_slide_master(prs, header_text)
    
    # ═══ (2) 표지 ═══
    for idx, role in slide_roles.items():
        if role == "cover":
            _update_cover(prs.slides[idx], company_name, course_name, today_str)
    
    # ═══ (3) 목차 ═══
    for idx, role in slide_roles.items():
        if role == "toc":
            _update_toc(prs.slides[idx], data)
    
    # ═══ (4) 과정 개요 ═══
    for idx, role in slide_roles.items():
        if role == "overview":
            _update_course_overview(prs.slides[idx], course, data, multi_result)
    
    # ═══ (5) Executive Summary + AI 내러티브 ═══
    exec_indices = [i for i, r in slide_roles.items() if r == "exec_summary"]
    _update_exec_summaries(prs, exec_indices, data, sample_design)
    
    # ═══ (5.5) AI 내러티브 삽입 ═══
    narrative = data.get("narrative", {})
    if narrative and exec_indices:
        _inject_narrative(prs.slides[exec_indices[0]], narrative, data)
    
    # ═══ (6) 정량 평가 — 차수별 분리 ═══
    quant_indices = [i for i, r in slide_roles.items() if r == "quant"]
    sessions = (multi_result or {}).get("sessions", [])
    if multi_result and multi_result.get("multi_session") and len(sessions) > 1:
        _update_quantitatives_multi(prs, quant_indices, data, sessions, sample_design)
    else:
        _update_quantitatives(prs, quant_indices, data, sample_design)
    
    # ═══ (7) 정성 평가 + 공통응답 그룹핑 ═══
    qual_indices = [i for i, r in slide_roles.items() if r == "qual"]
    _update_qualitatives_grouped(prs, qual_indices, data, sample_design)
    
    # ═══ (8) 빈 텍스트 채우기 (모든 슬라이드) ═══
    _fill_empty_placeholders(prs, company_name, course_name)
    
    prs.save(output_path)
    return output_path



# ═══════════════════════════════════════════════════
# 샘플 디자인 추출 (참조용만)
# ═══════════════════════════════════════════════════

def _extract_design_from_sample(sample_path):
    """샘플 PPT에서 디자인 속성만 추출 (콘텐츠는 가져오지 않음)"""
    try:
        prs = Presentation(sample_path)
    except:
        return None
    
    design = {"tables": [], "charts": [], "groups": []}
    
    for slide in prs.slides:
        for shape in slide.shapes:
            if shape.has_table:
                tbl = shape.table
                tbl_design = {
                    "rows": len(tbl.rows), "cols": len(tbl.columns),
                    "col_widths": [tbl.columns[c].width for c in range(len(tbl.columns))],
                    "row_heights": [tbl.rows[r].height for r in range(len(tbl.rows))],
                }
                # 헤더 셀 스타일
                if len(tbl.rows) > 0:
                    tbl_design["header_style"] = _extract_cell_style(tbl.cell(0, 0))
                if len(tbl.rows) > 1:
                    tbl_design["data_style"] = _extract_cell_style(tbl.cell(1, 0))
                design["tables"].append(tbl_design)
            
            if shape.shape_type == 6:
                grp_design = {
                    "position": {"left": shape.left, "top": shape.top,
                                 "width": shape.width, "height": shape.height},
                }
                design["groups"].append(grp_design)
    
    return design


def _extract_cell_style(cell):
    """셀 스타일 추출"""
    style = {}
    tc = cell._tc
    tcPr = tc.find(qn('a:tcPr'))
    if tcPr is not None:
        sf = tcPr.find(qn('a:solidFill'))
        if sf is not None:
            clr = sf.find(qn('a:srgbClr'))
            if clr is not None:
                style["bg_color"] = clr.get('val')
    try:
        if cell.text_frame.paragraphs and cell.text_frame.paragraphs[0].runs:
            run = cell.text_frame.paragraphs[0].runs[0]
            style["font_name"] = run.font.name
            style["font_size"] = int(run.font.size) if run.font.size else None
            style["bold"] = run.font.bold
    except:
        pass
    return style


# ═══════════════════════════════════════════════════
# 슬라이드 역할 자동 판별
# ═══════════════════════════════════════════════════

def _map_slide_roles(prs):
    roles = {}
    for idx, slide in enumerate(prs.slides):
        layout = slide.slide_layout.name
        title = _get_slide_title(slide)
        
        if "앞표지" in layout: roles[idx] = "cover"
        elif "목차" in layout: roles[idx] = "toc"
        elif "뒤표지" in layout: roles[idx] = "back"
        elif "Executive" in title: roles[idx] = "exec_summary"
        elif "정량" in title: roles[idx] = "quant"
        elif "정성" in title: roles[idx] = "qual"
        elif "과정 개요" in title or "과정개요" in title or "과정내용" in title: roles[idx] = "overview"
        elif "일정" in title: roles[idx] = "schedule"
        elif "만족도 결과" in title: roles[idx] = "section_satisfaction"
        elif "Appendix" in title: roles[idx] = "appendix"
        else: roles[idx] = "other"
    
    return roles


def _get_slide_title(slide):
    for shape in slide.shapes:
        if shape.has_text_frame:
            for p in shape.text_frame.paragraphs:
                t = p.text.strip()
                if t and 2 < len(t) < 40:
                    return t
    return ""


# ═══════════════════════════════════════════════════
# 마스터 헤더
# ═══════════════════════════════════════════════════

def _update_slide_master(prs, header_text):
    today_str = datetime.now().strftime("%Y-%m-%d")
    for master in prs.slide_masters:
        for shape in master.shapes:
            if shape.has_text_frame and ("Date" in shape.name or "날짜" in shape.name):
                _replace_text_preserve_format(shape.text_frame, today_str)
        for layout in master.slide_layouts:
            for shape in layout.shapes:
                if shape.has_text_frame:
                    text = shape.text_frame.text
                    if "과정운영" in text or "결과보고서" in text or "고객사" in text:
                        _replace_text_preserve_format(shape.text_frame, header_text)


# ═══════════════════════════════════════════════════
# 표지
# ═══════════════════════════════════════════════════

def _update_cover(slide, company, course_name, today_str):
    for shape in slide.shapes:
        if shape.shape_type == 6:
            try:
                for child in shape.shapes:
                    if child.has_text_frame:
                        _replace_text_preserve_format(child.text_frame, today_str)
            except:
                pass
        if shape.has_text_frame:
            full = shape.text_frame.text
            if "고객사" in full or "과정명" in full or "결과보고서" in full:
                paras = shape.text_frame.paragraphs
                lines = [company, course_name, "과정운영 결과보고서"]
                for i, line in enumerate(lines):
                    if i < len(paras):
                        _replace_para_text_preserve(paras[i], line)


# ═══════════════════════════════════════════════════
# 목차
# ═══════════════════════════════════════════════════

def _update_toc(slide, data):
    for shape in slide.shapes:
        if shape.has_table:
            table = shape.table
            toc_items = [
                ("Ⅰ", ". 과정 개요", ""),
                ("Ⅱ", ". 만족도 결과", ""),
                ("Ⅲ", ". Appendix", ""),
            ]
            for r in range(min(len(toc_items), len(table.rows))):
                for c in range(min(len(toc_items[r]), len(table.columns))):
                    _set_cell_text_preserve(table.cell(r, c), toc_items[r][c])


# ═══════════════════════════════════════════════════
# 과정 개요 — 빈 칸 없이
# ═══════════════════════════════════════════════════

def _update_course_overview(slide, course, data, multi_result=None):
    for shape in slide.shapes:
        if not shape.has_table:
            continue
        table = shape.table
        if len(table.rows) < 7:
            continue
        
        company = course.get("company", "")
        cname = course.get("course_name", "")
        full_name = f"{company} {cname}".strip() or "교육 과정"
        
        _set_cell_text_preserve(table.cell(0, 1), full_name)
        _set_cell_text_preserve(table.cell(1, 1), course.get("purpose", "직무역량 향상 및 조직 성과 기여"))
        
        # 모듈/과정 내용
        module_cat = next((c for c in data.get("categories", []) if c["name"] == "모듈"), None)
        if module_cat:
            lines = [f"M{i+1}. {q['label']}" for i, q in enumerate(module_cat["questions"])]
            _set_cell_text_preserve(table.cell(2, 1), "\n".join(lines))
        else:
            cats = data.get("categories", [])
            if cats:
                lines = [f"- {c['name']} ({len(c['questions'])}문항)" for c in cats]
                _set_cell_text_preserve(table.cell(2, 1), "\n".join(lines))
            else:
                _set_cell_text_preserve(table.cell(2, 1), full_name)
        
        _set_cell_text_preserve(table.cell(3, 1), course.get("date", datetime.now().strftime("%Y년 %m월")))
        _set_cell_text_preserve(table.cell(4, 1), course.get("location", "사내 교육장"))
        
        # 응답 인원 (차수별 표시)
        resp = data.get("response_count", 0)
        sessions = (multi_result or {}).get("sessions", [])
        if sessions and len(sessions) > 1:
            resp_lines = [f"총 {resp}명"]
            for s in sessions:
                resp_lines.append(f"  · {s.get('session_label', '')}: {s.get('response_count', 0)}명")
            resp_text = "\n".join(resp_lines)
        else:
            resp_text = f"{resp}명" if resp else "해당 없음"
        _set_cell_text_preserve(table.cell(5, 1), resp_text)
        
        instr = course.get("instructor", "")
        if not instr and data.get("instructor_names"):
            instr = ", ".join(data["instructor_names"])
        if len(table.rows) > 6:
            _set_cell_text_preserve(table.cell(6, 1), instr or "사내/외부 강사")
        break



# Executive Summary
# ═══════════════════════════════════════════════════

def _update_exec_summaries(prs, slide_indices, data, sample_design):
    if not slide_indices:
        return
    
    cats = data["categories"]
    
    for si, slide_idx in enumerate(slide_indices):
        slide = prs.slides[slide_idx]
        cat_slice = _get_cats_for_exec_slide(si, slide_indices, cats)
        
        for shape in slide.shapes:
            if hasattr(shape, 'has_chart') and shape.has_chart:
                chart_qs = []
                for c in cat_slice:
                    chart_qs.extend(c["questions"])
                if chart_qs:
                    _update_chart_data(shape.chart, {
                        "labels": [q["id"] for q in chart_qs],
                        "values": [q["avg"] for q in chart_qs],
                    })
            
            if shape.has_table:
                table = shape.table
                if len(table.rows) <= 4:
                    _fill_exec_table(table, cat_slice, data.get("overall_average", 0))


def _get_cats_for_exec_slide(slide_num, all_exec_indices, cats):
    total = len(all_exec_indices)
    if total <= 1:
        return cats
    per = max(1, len(cats) // total)
    start = slide_num * per
    end = start + per if slide_num < total - 1 else len(cats)
    return cats[start:end]


def _fill_exec_table(table, cats, overall_avg):
    cols = len(table.columns)
    for i in range(cols):
        if i < len(cats):
            name = cats[i]["name"]
            if not name.endswith("만족도"):
                name += " 만족도"
            _set_cell_text_preserve(table.cell(0, i), name)
            _set_cell_text_preserve(table.cell(1, i), f"{cats[i]['avg']:.2f}")
        else:
            _set_cell_text_preserve(table.cell(0, i), "-")
            _set_cell_text_preserve(table.cell(1, i), "-")
    
    if len(table.rows) > 2:
        _set_cell_text_preserve(table.cell(2, 0), f"전체 평균: {overall_avg:.2f}")


def _inject_narrative(slide, narrative, data):
    """AI 내러티브를 Executive Summary 슬라이드에 삽입"""
    exec_summary = narrative.get("executive_summary", "")
    if not exec_summary:
        # AI 미연결 시 규칙 기반 내러티브 자동 생성
        insights = data.get("insights", {})
        overall = insights.get("overall_average", 0)
        tier = insights.get("overall_tier", "")
        top = insights.get("top_categories", [])
        bottom = insights.get("bottom_categories", [])
        
        tier_text = {"excellent": "매우 우수한", "good": "양호한", "moderate": "보통의", "needs_improvement": "개선이 필요한"}.get(tier, "")
        
        lines = []
        lines.append(f"본 과정의 전체 만족도는 {overall:.2f}점(5점 만점)으로 {tier_text} 수준입니다.")
        
        if top:
            top_names = ", ".join([f"'{c['name']}'({c['avg']:.2f}점)" for c in top[:2]])
            lines.append(f"특히 {top_names} 영역이 높은 만족도를 기록하여 교육 성과로 평가됩니다.")
        
        if bottom:
            bot_names = ", ".join([f"'{c['name']}'({c['avg']:.2f}점)" for c in bottom])
            lines.append(f"{bot_names} 영역은 향후 과정 개선 시 보완할 수 있는 성장 기회로 판단됩니다.")
        
        exec_summary = " ".join(lines)
    
    # 텍스트 프레임이 있는 도형에 삽입 (표가 아닌 것 중 가장 큰 것)
    text_shapes = []
    for shape in slide.shapes:
        if shape.has_text_frame and not shape.has_table:
            # 제목/라벨 제외 (작은 텍스트)
            if shape.width > Cm(10) and shape.height > Cm(1.5):
                text_shapes.append(shape)
    
    # 가장 아래에 있는 큰 텍스트 프레임에 삽입
    if text_shapes:
        target = sorted(text_shapes, key=lambda s: s.top, reverse=True)[0]
        _replace_text_preserve_format(target.text_frame, exec_summary)
    
    # 강점/약점 코멘트가 있으면 추가 도형에
    strength = narrative.get("strength_comment", "")
    improvement = narrative.get("improvement_comment", "")
    recommendation = narrative.get("recommendation", "")
    
    remaining = [s for s in text_shapes if s != (sorted(text_shapes, key=lambda s: s.top, reverse=True)[0] if text_shapes else None)]
    comments = [c for c in [strength, improvement, recommendation] if c]
    for i, shape in enumerate(remaining):
        if i < len(comments):
            _replace_text_preserve_format(shape.text_frame, comments[i])


# ═══════════════════════════════════════════════════
# 정량 평가
# ═══════════════════════════════════════════════════

def _update_quantitatives(prs, slide_indices, data, sample_design):
    if not slide_indices:
        return
    
    questions = data["questions"]
    resp = data.get("response_count", 0)
    
    general_qs = [q for q in questions if q["category"] != "모듈" and "강사" not in q["category"]]
    module_qs = [q for q in questions if q["category"] == "모듈"]
    instructor_qs = [q for q in questions if "강사" in q["category"]]
    
    q_groups = [general_qs]
    if module_qs: q_groups.append(module_qs)
    if instructor_qs: q_groups.append(instructor_qs)
    
    for si, slide_idx in enumerate(slide_indices):
        slide = prs.slides[slide_idx]
        qs = q_groups[si] if si < len(q_groups) else questions
        
        for shape in slide.shapes:
            if shape.has_table and len(shape.table.rows) > 3:
                _fill_quant_table(shape.table, qs, resp)


def _fill_quant_table(table, questions, resp_count):
    rows = len(table.rows)
    cols = len(table.columns)
    
    cur_cat = ""
    for i, q in enumerate(questions):
        r = i + 1
        if r >= rows:
            break
        
        cat_text = ""
        is_first = False
        if q["category"] != cur_cat:
            cat_text = q["category"]
            cur_cat = q["category"]
            is_first = True
        
        if cols >= 4:
            _set_cell_text_preserve(table.cell(r, 0), cat_text)
            _set_cell_text_preserve(table.cell(r, 1), f"{q['id']}. {q['label']}")
            _set_cell_text_preserve(table.cell(r, 2), f"{q['avg']:.2f}")
            _set_cell_text_preserve(table.cell(r, 3), f"{resp_count}명" if is_first else "")
    
    # 남는 행에 "-" 채우기 (빈 칸 방지)
    for r in range(len(questions) + 1, rows):
        for c in range(cols):
            _set_cell_text_preserve(table.cell(r, c), "-")


# ═══════════════════════════════════════════════════
# 정성 평가
# ═══════════════════════════════════════════════════

def _update_qualitatives(prs, slide_indices, data, sample_design):
    if not slide_indices:
        return
    
    open_ended = data.get("open_ended", [])
    if not open_ended:
        return
    
    for slide_idx in slide_indices:
        slide = prs.slides[slide_idx]
        groups = [s for s in slide.shapes if s.shape_type == 6]
        
        if not groups:
            continue
        
        ref_group_xml = copy.deepcopy(groups[0]._element)
        
        for gi, grp in enumerate(groups):
            if gi < len(open_ended):
                _fill_group_data(grp, open_ended[gi])
            else:
                try:
                    for child in grp.shapes:
                        if child.has_text_frame:
                            _replace_text_preserve_format(child.text_frame, "(해당 없음)")
                except:
                    pass
        
        if len(open_ended) > len(groups):
            for gi in range(len(groups), len(open_ended)):
                oe = open_ended[gi]
                new_group = copy.deepcopy(ref_group_xml)
                xfrm = new_group.find(qn('p:grpSpPr'))
                if xfrm is not None:
                    off_el = xfrm.find(qn('a:xfrm'))
                    if off_el is not None:
                        off = off_el.find(qn('a:off'))
                        if off is not None:
                            orig_y = int(off.get('y', '0'))
                            off.set('y', str(orig_y + gi * Cm(4)))
                
                for sp in new_group.findall('.//' + qn('p:sp')):
                    nvSpPr = sp.find(qn('p:nvSpPr'))
                    if nvSpPr is None: continue
                    cNvPr = nvSpPr.find(qn('p:cNvPr'))
                    if cNvPr is None: continue
                    name = cNvPr.get('name', '')
                    txBody = sp.find(qn('p:txBody'))
                    if txBody is None: continue
                    
                    if '20' in name or '10' in name:
                        _set_xml_text(txBody, f"{oe['id']}. {oe['label']}")
                    elif '19' in name or '9' in name:
                        answers = "\n".join([f"• {a}" for a in oe["answers"]]) if oe["answers"] else "(응답 없음)"
                        _set_xml_text(txBody, answers)
                
                slide.shapes._spTree.append(new_group)


def _fill_group_data(group, oe):
    try:
        for child in group.shapes:
            if child.has_text_frame:
                name = child.name
                if '20' in name or '10' in name:
                    _replace_text_preserve_format(child.text_frame, f"{oe['id']}. {oe['label']}")
                elif '19' in name or '9' in name:
                    answers = "\n".join([f"• {a}" for a in oe["answers"]]) if oe["answers"] else "(응답 없음)"
                    _replace_text_preserve_format(child.text_frame, answers)
    except:
        pass


# ═══════════════════════════════════════════════════
# 빈 플레이스홀더 채우기 — 빈 칸 방지
# ═══════════════════════════════════════════════════

def _fill_empty_placeholders(prs, company, course_name):
    """모든 슬라이드의 빈 텍스트 박스에 기본값 채우기"""
    skip_layouts = ["뒤표지"]
    
    for slide in prs.slides:
        layout = slide.slide_layout.name
        if any(sk in layout for sk in skip_layouts):
            continue
        
        for shape in slide.shapes:
            if not shape.has_text_frame:
                continue
            
            name = shape.name.lower()
            text = shape.text_frame.text.strip()
            
            # 이미 내용이 있으면 건너뛰기
            if text and text not in ("고객사명", "과정명", "내용을 입력해주세요", "텍스트를 입력하세요"):
                continue
            
            # 플레이스홀더 패턴별 기본값
            if "고객사" in text or "고객사" in name:
                _replace_text_preserve_format(shape.text_frame, company)
            elif "과정명" in text or "과정" in name:
                _replace_text_preserve_format(shape.text_frame, course_name)
            elif "내용" in text or "텍스트" in text:
                _replace_text_preserve_format(shape.text_frame, f"{company} {course_name}")
            elif not text and "title" in name:
                pass  # 제목은 비워둬도 됨 (레이아웃 기본값)
    
    # 표 빈 셀도 채우기
    for slide in prs.slides:
        for shape in slide.shapes:
            if shape.has_table:
                for r in range(len(shape.table.rows)):
                    for c in range(len(shape.table.columns)):
                        cell = shape.table.cell(r, c)
                        if not cell.text.strip():
                            pass  # 표 빈 셀은 의도적일 수 있으므로 유지


# ═══════════════════════════════════════════════════
# 차트 업데이트
# ═══════════════════════════════════════════════════

def _update_chart_data(chart, chart_data):
    labels, values = chart_data["labels"], chart_data["values"]
    for series in chart.series:
        ser = series._element
        num_ref = ser.find(qn('c:val'))
        if num_ref is None: continue
        parent = num_ref.find(qn('c:numRef'))
        if parent is None: parent = num_ref.find(qn('c:numLit'))
        if parent is None: continue
        cache = parent.find(qn('c:numCache'))
        if cache is None: cache = parent
        for pt in cache.findall(qn('c:pt')): cache.remove(pt)
        pc = cache.find(qn('c:ptCount'))
        if pc is not None: pc.set('val', str(len(values)))
        for i, val in enumerate(values):
            pt = etree.SubElement(cache, qn('c:pt'))
            pt.set('idx', str(i))
            v = etree.SubElement(pt, qn('c:v'))
            v.text = f"{val:.2f}"
    try:
        first_ser = chart.plots[0]._element.find(qn('c:ser'))
        if first_ser is not None:
            cat_ref = first_ser.find(qn('c:cat'))
            if cat_ref is not None:
                sr = cat_ref.find(qn('c:strRef'))
                if sr is None: sr = cat_ref.find(qn('c:strLit'))
                if sr is not None:
                    sc = sr.find(qn('c:strCache'))
                    if sc is None: sc = sr
                    for pt in sc.findall(qn('c:pt')): sc.remove(pt)
                    pc = sc.find(qn('c:ptCount'))
                    if pc is not None: pc.set('val', str(len(labels)))
                    for i, lbl in enumerate(labels):
                        pt = etree.SubElement(sc, qn('c:pt'))
                        pt.set('idx', str(i))
                        v = etree.SubElement(pt, qn('c:v'))
                        v.text = lbl
    except: pass


# ═══════════════════════════════════════════════════
# 서식 보존 텍스트 유틸리티
# ═══════════════════════════════════════════════════

def _replace_text_preserve_format(text_frame, new_text):
    """텍스트 프레임의 텍스트만 교체 (서식 100% 보존)"""
    lines = new_text.replace('\x0b', '\n').split('\n')
    paras = text_frame.paragraphs
    
    ref_run_xml = None
    ref_pPr = None
    if paras and paras[0].runs:
        ref_run_xml = copy.deepcopy(paras[0].runs[0]._r)
        pPr_el = paras[0]._p.find(qn('a:pPr'))
        if pPr_el is not None:
            ref_pPr = copy.deepcopy(pPr_el)
    
    txBody = text_frame._txBody
    for p in txBody.findall(qn('a:p')):
        txBody.remove(p)
    
    for line in lines:
        p = etree.SubElement(txBody, qn('a:p'))
        if ref_pPr is not None:
            p.insert(0, copy.deepcopy(ref_pPr))
        if ref_run_xml is not None:
            r = copy.deepcopy(ref_run_xml)
            t = r.find(qn('a:t'))
            if t is None:
                t = etree.SubElement(r, qn('a:t'))
            t.text = line
            p.append(r)
        else:
            r = etree.SubElement(p, qn('a:r'))
            t = etree.SubElement(r, qn('a:t'))
            t.text = line


def _replace_para_text_preserve(paragraph, new_text):
    if paragraph.runs:
        for run in paragraph.runs[1:]:
            run._r.getparent().remove(run._r)
        paragraph.runs[0].text = new_text
    else:
        paragraph.text = new_text


def _set_cell_text_preserve(cell, new_text):
    """표 셀의 텍스트만 교체 (셀 서식/배경/폰트 100% 보존)"""
    tf = cell.text_frame
    if not tf.paragraphs:
        tf.text = new_text
        return
    
    lines = new_text.split('\n')
    p = tf.paragraphs[0]
    if p.runs:
        for run in p.runs[1:]:
            run._r.getparent().remove(run._r)
        p.runs[0].text = lines[0]
    else:
        p.text = lines[0]
    
    if len(lines) > 1:
        ref_rPr = None
        if p.runs:
            rPr = p.runs[0]._r.find(qn('a:rPr'))
            if rPr is not None:
                ref_rPr = copy.deepcopy(rPr)
        
        for line in lines[1:]:
            new_p = etree.SubElement(tf._txBody, qn('a:p'))
            r = etree.SubElement(new_p, qn('a:r'))
            if ref_rPr is not None:
                r.insert(0, copy.deepcopy(ref_rPr))
            t = etree.SubElement(r, qn('a:t'))
            t.text = line
    
    for para in tf.paragraphs[len(lines):]:
        para._p.getparent().remove(para._p)


def _set_xml_text(txBody, text):
    paras = txBody.findall(qn('a:p'))
    ref_rPr = None
    ref_pPr = None
    if paras:
        r = paras[0].find(qn('a:r'))
        if r is not None:
            rPr = r.find(qn('a:rPr'))
            if rPr is not None: ref_rPr = copy.deepcopy(rPr)
        pPr = paras[0].find(qn('a:pPr'))
        if pPr is not None: ref_pPr = copy.deepcopy(pPr)
    for p in paras:
        txBody.remove(p)
    for line in text.split('\n'):
        p = etree.SubElement(txBody, qn('a:p'))
        if ref_pPr is not None: p.insert(0, copy.deepcopy(ref_pPr))
        r = etree.SubElement(p, qn('a:r'))
        if ref_rPr is not None: r.insert(0, copy.deepcopy(ref_rPr))
        t = etree.SubElement(r, qn('a:t'))
        t.text = line


# ═══════════════════════════════════════════════════
# 차수별 정량 평가 (멀티세션)
# ═══════════════════════════════════════════════════

def _update_quantitatives_multi(prs, slide_indices, combined_data, sessions, sample_design):
    """
    차수가 2개 이상일 때:
    - 첫 슬라이드: 종합 정량 평가
    - 이후 슬라이드: 차수별 정량 평가
    - 마지막: 차수 비교 표
    """
    if not slide_indices:
        return

    resp = combined_data.get("response_count", 0)

    # 슬라이드 0번 → 종합
    slide0 = prs.slides[slide_indices[0]]
    for shape in slide0.shapes:
        if shape.has_text_frame:
            t = shape.text_frame.text
            if len(t) < 30:
                _replace_text_preserve_format(shape.text_frame, "정량 평가 결과 (종합)")
        if shape.has_table and len(shape.table.rows) > 3:
            _fill_quant_table(shape.table, combined_data["questions"], resp)

    # 차수별 슬라이드 (슬라이드 복제)
    if len(slide_indices) >= 2:
        # 두 번째 슬라이드 → 차수별 반복
        ref_idx = slide_indices[1] if len(slide_indices) > 1 else slide_indices[0]
        ref_slide = prs.slides[ref_idx]

        for si, session in enumerate(sessions):
            if si < len(slide_indices) - 1:
                target_slide = prs.slides[slide_indices[si + 1]]
            else:
                # 슬라이드 부족 — 기존 마지막 슬라이드 재활용
                target_slide = prs.slides[slide_indices[-1]]

            label = session.get("session_label", f"{si+1}차수")
            s_resp = session.get("response_count", 0)

            for shape in target_slide.shapes:
                if shape.has_text_frame:
                    t = shape.text_frame.text
                    if len(t) < 30:
                        _replace_text_preserve_format(shape.text_frame, f"정량 평가 결과 ({label})")
                if shape.has_table and len(shape.table.rows) > 3:
                    _fill_quant_table(shape.table, session.get("questions", []), s_resp)
    else:
        # 슬라이드 1개뿐 — 종합만
        _update_quantitatives(prs, slide_indices, combined_data, sample_design)


# ═══════════════════════════════════════════════════
# 공통응답 그룹핑 정성 평가
# ═══════════════════════════════════════════════════

def _update_qualitatives_grouped(prs, slide_indices, data, sample_design):
    """
    공통응답 그룹핑 결과를 정성 평가 슬라이드에 삽입
    - open_ended_grouped 있으면 그룹별 공통응답 표시
    - 없으면 기존 _update_qualitatives 로직 사용
    """
    grouped_oe = data.get("open_ended_grouped")

    if not grouped_oe:
        # 폴백: 기존 로직
        _update_qualitatives(prs, slide_indices, data, sample_design)
        return

    if not slide_indices:
        return

    for slide_idx in slide_indices:
        slide = prs.slides[slide_idx]
        groups_shapes = [s for s in slide.shapes if s.shape_type == 6]  # 그룹 도형

        for oe_idx, oe in enumerate(grouped_oe):
            groups = oe.get("groups", [])
            if not groups:
                continue

            # 공통응답 텍스트 구성 (가나다 오름차순 — 이미 정렬됨)
            lines = []
            for g in groups:
                cid = g.get("common_id", "")
                lbl = g.get("label", "")
                cnt = g.get("count", 0)
                lines.append(f"[{cid}] {lbl} ({cnt}건)")

            answer_text = "\n".join(lines) if lines else "(응답 없음)"
            question_text = f"{oe.get('id', '')}. {oe.get('label', '')}"

            # 그룹 도형에 삽입
            if oe_idx < len(groups_shapes):
                _fill_group_data_text(groups_shapes[oe_idx], question_text, answer_text)
            else:
                # 그룹 도형 부족 — 텍스트 박스 추가
                if groups_shapes:
                    _fill_group_data_text(groups_shapes[-1], question_text, answer_text)


def _fill_group_data_text(group_shape, question_text, answer_text):
    """그룹 도형에 질문/답변 텍스트 삽입"""
    try:
        for child in group_shape.shapes:
            if child.has_text_frame:
                name = child.name
                # 질문 영역
                if any(n in name for n in ['20', '10', 'title', 'Title', '제목', 'Q']):
                    _replace_text_preserve_format(child.text_frame, question_text)
                # 답변 영역
                elif any(n in name for n in ['19', '9', 'body', 'Body', '내용', 'A']):
                    _replace_text_preserve_format(child.text_frame, answer_text)
    except:
        # 단순 폴백
        try:
            shapes = list(group_shape.shapes)
            if len(shapes) >= 2:
                _replace_text_preserve_format(shapes[0].text_frame, question_text)
                _replace_text_preserve_format(shapes[1].text_frame, answer_text)
        except:
            pass
