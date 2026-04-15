"""코스모: 샘플 vs 출력 정밀 비교"""
import sys, glob, os
sys.stdout.reconfigure(encoding='utf-8')
from pptx import Presentation

# 코스모 샘플 / 출력
sample = next(f for f in glob.glob('samples/*.pptx') if '2026' in os.path.basename(f) and '코스모' in os.path.basename(f))
output = next(f for f in glob.glob('output/*.pptx') if '코스모' in os.path.basename(f) or '주니어' in os.path.basename(f))

print(f"샘플: {os.path.basename(sample)} ({os.path.getsize(sample)//1024}KB)")
print(f"출력: {os.path.basename(output)} ({os.path.getsize(output)//1024}KB)")

srs = Presentation(sample)
ors = Presentation(output)
print(f"슬라이드: 샘플={len(srs.slides)} 출력={len(ors.slides)}")

print("\n--- 구조 + 디자인 비교 ---")
for i in range(min(len(srs.slides), len(ors.slides))):
    ss = srs.slides[i]
    os_ = ors.slides[i]
    
    s_layout = ss.slide_layout.name
    o_layout = os_.slide_layout.name
    layout_ok = s_layout == o_layout
    
    s_title = ""
    o_title = ""
    for sh in ss.shapes:
        if sh.has_text_frame:
            for p in sh.text_frame.paragraphs:
                t = p.text.strip()
                if t and 2 < len(t) < 35 and not s_title: s_title = t
    for sh in os_.shapes:
        if sh.has_text_frame:
            for p in sh.text_frame.paragraphs:
                t = p.text.strip()
                if t and 2 < len(t) < 35 and not o_title: o_title = t
    
    s_shapes = len(list(ss.shapes))
    o_shapes = len(list(os_.shapes))
    
    s_tables = [(sh.table, sh.name) for sh in ss.shapes if sh.has_table]
    o_tables = [(sh.table, sh.name) for sh in os_.shapes if sh.has_table]
    
    status = "OK" if layout_ok and s_shapes == o_shapes else "!!"
    print(f"\nS{i+1:2d} [{status}] 레이아웃: {s_layout}")
    print(f"    샘플: {s_title} (도형 {s_shapes}개)")
    print(f"    출력: {o_title} (도형 {o_shapes}개)")
    
    # 표 비교
    for j in range(min(len(s_tables), len(o_tables))):
        st, sn = s_tables[j]
        ot, on = o_tables[j]
        s_rc = f"{len(st.rows)}x{len(st.columns)}"
        o_rc = f"{len(ot.rows)}x{len(ot.columns)}"
        size_ok = s_rc == o_rc
        
        # 첫 행 헤더
        s_h = [st.cell(0, c).text[:15] for c in range(len(st.columns))]
        o_h = [ot.cell(0, c).text[:15] for c in range(len(ot.columns))]
        
        tbl_status = "OK" if size_ok else "!!"
        print(f"    표 [{tbl_status}] 샘플={s_rc} 출력={o_rc}")
        if not size_ok or s_h != o_h:
            print(f"         헤더S: {s_h}")
            print(f"         헤더O: {o_h}")

print("\n--- 완료 ---")
