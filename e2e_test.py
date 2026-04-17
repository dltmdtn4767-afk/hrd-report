"""E2E 테스트: 엑셀 업로드 → PPT 생성 → 파일 검증"""
import requests
import os
import glob

BASE = "http://localhost:8080"

# 1. 엑셀 파일 찾기
xlsx_files = glob.glob(r"C:\Users\User\Desktop\업무폴더\*.xlsx")
print(f"=== 1. Excel files found: {len(xlsx_files)} ===")
for f in xlsx_files:
    print(f"   {os.path.basename(f)} ({os.path.getsize(f)} bytes)")

# 작은 파일 사용
test_file = min(xlsx_files, key=os.path.getsize)
print(f"\n   Using: {os.path.basename(test_file)}")

# 2. 업로드
print(f"\n=== 2. Uploading Excel ===")
with open(test_file, 'rb') as f:
    r = requests.post(f"{BASE}/api/upload", files={"file": f})
print(f"   Status: {r.status_code}")
if r.status_code != 200:
    print(f"   ERROR: {r.text[:500]}")
    exit(1)

data = r.json()
sid = data["session_id"]
summary = data["summary"]
print(f"   Session: {sid}")
print(f"   Course: {summary.get('course_name','?')}")
print(f"   Questions: {summary.get('total_questions',0)}")
print(f"   Open-ended: {summary.get('open_ended_count',0)}")
print(f"   Response count: {summary.get('response_count',0)}")

# 3. PPT 생성
print(f"\n=== 3. Generating PPT ===")
payload = {
    "slides": [
        {"type": "cover", "data": {"company": summary.get("company",""), "course": summary.get("course_name","")}},
        {"type": "toc", "data": {}},
        {"type": "section_quant", "data": {}},
        {"type": "overview", "data": {"rows": []}},
        {"type": "schedule", "data": {"rows": []}},
        {"type": "summary", "data": {}},
        {"type": "quant_chart", "groupId": 1, "data": {}},
        {"type": "section_qual", "data": {}},
        {"type": "qual_text", "data": {}},
        {"type": "back_cover", "data": {}}
    ],
    "quant_groups": [{"id": 1, "title": "전체 문항", "questions": []}],
    "qual_data": []
}

r = requests.post(f"{BASE}/api/build_ppt/{sid}", json=payload)
print(f"   Status: {r.status_code}")
print(f"   Content-Type: {r.headers.get('content-type','?')}")
print(f"   Content-Length: {len(r.content)} bytes")

if r.status_code != 200:
    print(f"   ERROR: {r.text[:500]}")
    exit(1)

# 4. 파일 검증
out_file = "e2e_test_output.pptx"
with open(out_file, "wb") as f:
    f.write(r.content)
print(f"   Saved: {out_file} ({os.path.getsize(out_file)} bytes)")

# PPTX 검증 (ZIP + python-pptx)
import zipfile
is_zip = zipfile.is_zipfile(out_file)
print(f"\n=== 4. Validation ===")
print(f"   Is valid ZIP: {is_zip}")

if is_zip:
    from pptx import Presentation
    p = Presentation(out_file)
    print(f"   Is valid PPTX: True")
    print(f"   Slide count: {len(p.slides)}")
    for i, sl in enumerate(p.slides):
        texts = []
        for sh in sl.shapes:
            if sh.has_text_frame:
                t = sh.text_frame.text.strip()[:40]
                if t:
                    texts.append(t)
        print(f"   Slide {i+1}: {texts[:2]}")
else:
    # 파일 첫 바이트 확인
    with open(out_file, 'rb') as f:
        header = f.read(8)
    print(f"   File header bytes: {header}")
    print(f"   NOT a valid PPTX file!")

# Content-Disposition 확인
cd = r.headers.get('content-disposition', '')
print(f"\n   Content-Disposition: {cd}")

print(f"\n=== DONE ===")
