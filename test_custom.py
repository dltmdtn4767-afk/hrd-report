import requests, os, glob, zipfile
BASE = "http://localhost:8080"

# 1. Upload
xlsx = min(glob.glob(r"C:\Users\User\Desktop\업무폴더\*.xlsx"), key=os.path.getsize)
with open(xlsx, "rb") as f:
    r = requests.post(f"{BASE}/api/upload", files={"file": f})
d = r.json()
sid = d["session_id"]
print("Upload:", r.status_code, sid)

# 2. PPT with custom_quant
payload = {
    "slides": [
        {"type": "cover", "data": {"company": "Test Corp", "course": "Test Course"}},
        {"type": "summary", "data": {}},
        {"type": "custom_quant", "customSlideId": 1, "data": {
            "title": "Grouped Analysis",
            "chartType": "bar",
            "tableStyle": "B",
            "tablePos": "below",
            "groups": [
                {"name": "Group A", "avg": 4.25, "color": "#36A86F", "qIds": [],
                 "questions": [{"id":"Q1","label":"Question 1","avg":4.5,"count":50},
                               {"id":"Q2","label":"Question 2","avg":4.0,"count":50}]},
                {"name": "Group B", "avg": 3.80, "color": "#4A90D9", "qIds": [],
                 "questions": [{"id":"Q3","label":"Question 3","avg":3.8,"count":50}]}
            ]
        }},
        {"type": "qual_text", "data": {}},
        {"type": "back_cover", "data": {}}
    ],
    "quant_groups": [],
    "qual_data": []
}
r = requests.post(f"{BASE}/api/build_ppt/{sid}", json=payload)
print("PPT:", r.status_code, "size:", len(r.content))

if r.status_code == 200:
    with open("test_custom.pptx", "wb") as f:
        f.write(r.content)
    print("Valid PPTX:", zipfile.is_zipfile("test_custom.pptx"))
    from pptx import Presentation
    p = Presentation("test_custom.pptx")
    print("Slides:", len(p.slides))
    for i, sl in enumerate(p.slides):
        texts = [sh.text_frame.text.strip()[:40] for sh in sl.shapes if sh.has_text_frame and sh.text_frame.text.strip()]
        print(f"  Slide {i+1}: {texts[:2]}")
else:
    print("ERROR:", r.text[:300])
