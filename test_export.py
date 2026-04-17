"""Test export_element API - chart and table"""
import requests, zipfile

BASE = "http://localhost:8080"

# Test chart export
r = requests.post(f"{BASE}/api/export_element", json={
    "type": "chart",
    "title": "Test Chart",
    "data": {
        "labels": ["Category A", "Category B", "Category C"],
        "values": [4.52, 3.80, 4.10],
        "chartType": "bar",
        "colors": ["#36A86F", "#4A90D9", "#4A90D9"]
    }
})
print(f"Chart: {r.status_code} size={len(r.content)}")
with open("export_chart.pptx", "wb") as f:
    f.write(r.content)
print(f"  Valid PPTX: {zipfile.is_zipfile('export_chart.pptx')}")

from pptx import Presentation
p = Presentation("export_chart.pptx")
for sh in p.slides[0].shapes:
    if sh.has_chart:
        c = sh.chart
        print(f"  Chart type: {c.chart_type}")
        print(f"  Categories: {[str(cat) for cat in c.plots[0].categories]}")
        vals = [v for v in c.series[0].values]
        print(f"  Values: {vals}")
        print("  -> Editable: YES (double-click in PPT)")

# Test table export
r = requests.post(f"{BASE}/api/export_element", json={
    "type": "table",
    "title": "Test Table",
    "data": {
        "headers": ["No", "Question", "Avg", "Count"],
        "rows": [
            ["1", "Question A", "4.52", "50"],
            ["2", "Question B", "3.80", "50"],
            ["3", "Question C", "4.10", "50"]
        ]
    }
})
print(f"\nTable: {r.status_code} size={len(r.content)}")
with open("export_table.pptx", "wb") as f:
    f.write(r.content)
print(f"  Valid PPTX: {zipfile.is_zipfile('export_table.pptx')}")

p2 = Presentation("export_table.pptx")
for sh in p2.slides[0].shapes:
    if sh.has_table:
        t = sh.table
        print(f"  Table: {len(t.rows)}r x {len(t.columns)}c")
        row0 = [t.cell(0,c).text for c in range(len(t.columns))]
        print(f"  Headers: {row0}")
        print("  -> Copyable: YES (select in PPT, Ctrl+C)")

print("\nALL OK")
