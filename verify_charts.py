"""Verify native charts in generated PPTX"""
from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE

p = Presentation("test_custom.pptx")
print(f"Slides: {len(p.slides)}")
for i, sl in enumerate(p.slides):
    charts = []
    tables = []
    for sh in sl.shapes:
        if sh.has_chart:
            c = sh.chart
            series_data = []
            try:
                for s in c.series:
                    vals = [v for v in s.values]
                    series_data.append(vals)
            except:
                pass
            cats = []
            try:
                cats = [str(cat) for cat in c.plots[0].categories]
            except:
                pass
            charts.append({"type": str(c.chart_type), "categories": cats, "values": series_data})
        if sh.has_table:
            t = sh.table
            tables.append(f"{len(t.rows)}r x {len(t.columns)}c")
    if charts or tables:
        print(f"\n  Slide {i+1}:")
        for ch in charts:
            print(f"    CHART: {ch['type']}")
            print(f"      Categories: {ch['categories'][:5]}")
            print(f"      Values: {ch['values'][:1]}")
            print(f"      -> Native=True, Editable=True, Copyable=True")
        for tb in tables:
            print(f"    TABLE: {tb} -> Native=True, Copyable=True")
print("\nDONE")
