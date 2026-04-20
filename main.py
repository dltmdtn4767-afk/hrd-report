"""
HRD 만족도 결과보고서 웹앱 — FastAPI 메인 서버 v3
멀티시트(차수별) + 공통응답 그룹핑 지원
"""
import os
import sys
import json
import uuid
import shutil
from datetime import datetime
from pathlib import Path

from fastapi import FastAPI, UploadFile, File, Form, HTTPException
from fastapi.staticfiles import StaticFiles
from fastapi.responses import FileResponse, HTMLResponse, JSONResponse
from fastapi.middleware.cors import CORSMiddleware

# 프로젝트 루트
BASE_DIR = Path(__file__).parent
sys.path.insert(0, str(BASE_DIR))

from modules.config_manager import load_config
from modules.data_loader import load_from_excel, load_all_sheets, get_available_sheets
from modules.analyzer import analyze_data, generate_narrative_prompt, generate_qualitative_prompt
from modules.response_grouper import process_all_open_ended
from modules.report_generator import generate_report, find_template
from modules.ai_engine import AIEngine
from modules.sample_analyzer import SampleAnalyzer
from modules.preview_engine import generate_preview

app = FastAPI(title="HRD 만족도 결과보고서 생성기", version="3.0")

app.add_middleware(
    CORSMiddleware,
    allow_origins=["*"],
    allow_methods=["*"],
    allow_headers=["*"],
)

app.mount("/static", StaticFiles(directory=str(BASE_DIR / "static")), name="static")

# 글로벌
config = load_config()
ai_engine = AIEngine(config)
sample_analyzer = SampleAnalyzer(config)

# 세션 저장소
sessions = {}


@app.get("/", response_class=HTMLResponse)
async def index():
    return (BASE_DIR / "static" / "index.html").read_text(encoding="utf-8")


@app.post("/api/upload")
async def upload_excel(file: UploadFile = File(...), sheet: str = Form(None)):
    """
    엑셀 업로드 → 전체 시트 분석 → AI 추론 → 구성안 반환
    - sheet=None : 모든 시트 자동 분석 (권장)
    - sheet=특정값 : 해당 시트만 분석
    """
    session_id = str(uuid.uuid4())[:8]
    upload_dir = BASE_DIR / "uploads" / session_id
    upload_dir.mkdir(parents=True, exist_ok=True)
    
    file_path = upload_dir / file.filename
    with open(file_path, "wb") as f:
        shutil.copyfileobj(file.file, f)
    
    # ═══ 멀티시트 로딩 ═══
    try:
        if sheet:
            # 특정 시트만
            raw_data = load_from_excel(str(file_path), sheet)
            multi_result = {
                "course_info": raw_data["course_info"],
                "sessions": [raw_data],
                "combined": raw_data,
                "multi_session": False,
            }
        else:
            # 전체 시트 자동 분석
            multi_result = load_all_sheets(str(file_path))
    except Exception as e:
        raise HTTPException(400, f"엑셀 파싱 실패: {str(e)}")
    
    # 분석 대상: 종합 데이터
    data = multi_result["combined"]
    data = analyze_data(data)
    multi_result["combined"] = data
    
    # 각 차수도 분석
    for s in multi_result["sessions"]:
        try:
            multi_result["sessions"][multi_result["sessions"].index(s)] = analyze_data(s)
        except:
            pass
    
    # 데이터 요약
    course = data["course_info"]
    cats = data["categories"]
    insights = data.get("insights", {})
    sessions_info = [
        {"label": s.get("session_label", s.get("sheet_name", "")),
         "response_count": s.get("response_count", 0),
         "overall_average": s.get("overall_average", 0)}
        for s in multi_result["sessions"]
    ]
    
    summary = {
        "company": course.get("company", ""),
        "course_name": course.get("course_name", ""),
        "total_questions": len(data["questions"]),
        "categories": len(cats),
        "category_names": [c["name"] for c in cats],
        "has_modules": any(c["name"] == "모듈" for c in cats),
        "num_instructors": len([c for c in cats if "강사" in c["name"]]),
        "open_ended_count": len(data.get("open_ended", [])),
        "response_count": data.get("response_count", 0),
        "overall_average": round(data.get("overall_average", 0), 2),
        "overall_tier": insights.get("overall_tier", ""),
        "consistency": insights.get("consistency", ""),
        "session_count": len(multi_result["sessions"]),
        "multi_session": multi_result["multi_session"],
        "sessions_info": sessions_info,
    }
    
    # AI 추론 + 샘플 매칭
    sample_patterns = sample_analyzer.get_patterns()
    ai_result = await ai_engine.infer_structure(summary, sample_patterns)
    
    matched_sample = sample_analyzer.find_best_match(summary)
    if matched_sample:
        ai_result["matched_sample"] = matched_sample["name"]
        ai_result["matched_slide_count"] = matched_sample["slide_count"]
    
    # 세션 저장
    sessions[session_id] = {
        "file_path": str(file_path),
        "sheet": sheet,
        "data": data,
        "multi_result": multi_result,
        "summary": summary,
        "ai_result": ai_result,
        "matched_sample": matched_sample,
        "created_at": datetime.now().isoformat(),
    }
    
    # 프론트 대시보드용 직렬화 데이터
    def _ser_session(s):
        return {
            "session_label": s.get("session_label", s.get("sheet_name", "")),
            "sheet_name": s.get("sheet_name", ""),
            "response_count": s.get("response_count", 0),
            "overall_average": s.get("overall_average", 0),
            "categories": [
                {"name": c["name"], "avg": c["avg"],
                 "per_session": c.get("per_session", {}),
                 "questions": [{"id": q["id"], "label": q["label"],
                                "avg": q["avg"], "count": q.get("count", 0)}
                               for q in c["questions"]]
                } for c in s.get("categories", [])
            ],
            "questions": [{"id": q["id"], "label": q["label"],
                           "avg": q["avg"], "count": q.get("count", 0),
                           "category": q["category"]}
                          for q in s.get("questions", [])],
            "open_ended": [{"id": o["id"], "label": o["label"],
                            "answers": o.get("answers", [])}
                           for o in s.get("open_ended", [])],
        }

    serialized_multi = {
        "multi_session": multi_result["multi_session"],
        "sessions": [_ser_session(s) for s in multi_result["sessions"]],
        "combined": _ser_session(multi_result["combined"]),
    }

    return {
        "session_id": session_id,
        "summary": summary,
        "ai_result": ai_result,
        "multi_result": serialized_multi,
    }


@app.get("/api/analyze_qual/{session_id}")
async def analyze_qual(session_id: str):
    """주관식 공통응답 그룹핑 (대시보드용)"""
    if session_id not in sessions:
        raise HTTPException(404, "세션 없음")
    sess = sessions[session_id]
    data = sess["data"]
    multi_result = sess.get("multi_result", {})
    combined = multi_result.get("combined", data)

    try:
        grouped = await process_all_open_ended(
            combined.get("open_ended", []), ai_engine
        )
        # 각 차수 주관식도 그룹핑
        sessions_qual = []
        for s in multi_result.get("sessions", []):
            sg = await process_all_open_ended(s.get("open_ended", []), ai_engine)
            sessions_qual.append({
                "session_label": s.get("session_label", s.get("sheet_name", "")),
                "open_ended_grouped": sg,
            })
        return {
            "open_ended_grouped": grouped,
            "sessions_qual": sessions_qual,
        }
    except Exception as e:
        import traceback; traceback.print_exc()
        return {"open_ended_grouped": [], "error": str(e)}


@app.get("/api/rawdata/{session_id}")
async def get_rawdata(session_id: str):
    """문항별 개별 응답(로우데이터) 반환 — Excel 상세 시트용"""
    if session_id not in sessions:
        raise HTTPException(404, "세션 없음")
    sess = sessions[session_id]
    file_path = sess["file_path"]
    multi_result = sess.get("multi_result", {})

    try:
        import openpyxl
        wb = openpyxl.load_workbook(file_path, data_only=True)

        raw_by_sheet = {}
        for sheet_name in wb.sheetnames:
            ws = wb[sheet_name]
            # 각 시트의 헤더 행 찾기
            from modules.data_loader import _find_header_row, parse_header, OPEN_ENDED_KEYWORDS
            header_row = _find_header_row(ws)
            headers = {}
            for col_idx in range(1, ws.max_column + 1):
                val = ws.cell(row=header_row, column=col_idx).value
                if not val:
                    continue
                skip = ['타임스탬프','timestamp','응답자','이름','이메일','email','제출일']
                if any(k in str(val).lower() for k in skip):
                    continue
                parsed = parse_header(str(val).strip())
                if parsed:
                    headers[col_idx] = {"id": parsed["id"], "label": parsed["label"],
                                        "is_open_ended": parsed["is_open_ended"], "responses": []}

            # 데이터 수집
            for row in range(header_row + 1, ws.max_row + 1):
                row_empty = all(
                    ws.cell(row=row, column=c).value is None
                    for c in headers
                )
                if row_empty:
                    continue
                for col_idx, q in headers.items():
                    val = ws.cell(row=row, column=col_idx).value
                    if val is not None:
                        q["responses"].append(val)

            raw_by_sheet[sheet_name] = list(headers.values())

        wb.close()
        return {"raw_by_sheet": raw_by_sheet}

    except Exception as e:
        import traceback; traceback.print_exc()
        raise HTTPException(500, f"로우데이터 로드 실패: {str(e)}")


@app.post("/api/generate/{session_id}")

async def generate(session_id: str):
    """PPT 생성: 차수별 슬라이드 + 종합 + 공통응답"""
    if session_id not in sessions:
        raise HTTPException(404, "세션 없음")
    
    session = sessions[session_id]
    data = session["data"]
    multi_result = session.get("multi_result", {})
    
    try:
        # ═══ (1) 공통응답 그룹핑 (AI 실패해도 규칙 기반으로 진행) ═══
        try:
            if data.get("open_ended"):
                grouped_oe = await process_all_open_ended(data["open_ended"], ai_engine)
                data["open_ended_grouped"] = grouped_oe
                for s in multi_result.get("sessions", []):
                    if s.get("open_ended"):
                        s["open_ended_grouped"] = await process_all_open_ended(s["open_ended"], ai_engine)
        except Exception as e_qual:
            print(f"[WARN] 정성 그룹핑 실패(폴백 사용): {e_qual}")
        
        # ═══ (2) AI 내러티브 — 실패해도 계속 ═══
        narrative = {}
        try:
            narrative_prompt = generate_narrative_prompt(data)
            narrative = await ai_engine.generate_narrative(narrative_prompt)
            data["narrative"] = narrative
        except Exception as e_narr:
            print(f"[WARN] AI 내러티브 실패: {e_narr}")
            data["narrative"] = {}
        
        # ═══ (3) 주관식 AI 테마 — 실패해도 계속 ═══
        try:
            if data.get("open_ended"):
                qual_prompt = generate_qualitative_prompt(data["open_ended"])
                qual_summary = await ai_engine.summarize_qualitative(qual_prompt)
                data["qualitative_summary"] = qual_summary
        except Exception as e_qual2:
            print(f"[WARN] 주관식 AI 요약 실패: {e_qual2}")
        
        # ═══ (4) PPT 생성 ═══
        matched = session.get("matched_sample")
        output_path = generate_report(
            data, config,
            sample_pattern=matched,
            multi_result=multi_result,
        )
        session["output_path"] = output_path
        
        preview = generate_preview(output_path)
        session["preview"] = preview
        
        # ═══ (5) AI 검토 — 실패해도 계속 ═══
        review = {}
        try:
            review = await ai_engine.review_output(data, preview)
            session["review"] = review
        except Exception as e_rev:
            print(f"[WARN] AI 검토 실패: {e_rev}")
        
        return {
            "success": True,
            "output_file": os.path.basename(output_path),
            "slide_count": len(preview),
            "preview": preview,
            "review": review,
            "narrative": narrative,
            "session_count": len(multi_result.get("sessions", [])),
            "multi_session": multi_result.get("multi_session", False),
        }

    except Exception as e:
        import traceback
        traceback.print_exc()
        return {"success": False, "error": str(e)}


@app.post("/api/build_ppt/{session_id}")
async def build_ppt(session_id: str, payload: dict):
    """보고서 빌더 → 커스텀 PPT 생성"""
    if session_id not in sessions:
        raise HTTPException(404, "세션 없음")
    sess = sessions[session_id]

    try:
        from modules.builder_generator import build_custom_ppt
        slides    = payload.get("slides", [])
        q_groups  = payload.get("quant_groups", [])
        qual_data = payload.get("qual_data", [])
        data      = sess.get("data", {})
        # 세션에 업로드 템플릿 있으면 우선 사용
        tpl_path  = sess.get("template_path", None)

        output_path = build_custom_ppt(slides, q_groups, qual_data, data, config,
                                       template_path=tpl_path)
        filename = os.path.basename(output_path)
        # 한글 파일명 → RFC 5987 인코딩
        from urllib.parse import quote
        safe_fn = quote(filename)
        return FileResponse(
            output_path,
            media_type="application/vnd.openxmlformats-officedocument.presentationml.presentation",
            headers={"Content-Disposition": f"attachment; filename*=UTF-8''{safe_fn}"},
        )
    except Exception as e:
        import traceback; traceback.print_exc()
        raise HTTPException(500, f"빌더 PPT 생성 실패: {e}")


@app.post("/api/export_element")
async def export_element(payload: dict):
    """차트 또는 표 하나만 들어있는 PPTX 생성 → 복사 붙여넣기용
    payload: { type: 'chart'|'table', title: str, data: {...} }
    """
    from pptx import Presentation as Prs
    from pptx.util import Inches, Pt, Emu, Cm
    from pptx.dml.color import RGBColor
    from pptx.chart.data import ChartData
    from pptx.enum.chart import XL_CHART_TYPE
    from pptx.enum.text import PP_ALIGN
    from urllib.parse import quote
    from modules.ppt_constants import (
        FONT_NAME, BRAND_COLORS,
        CHART_SIZE, TABLE_SIZE,
        CHART_CONFIG
    )
    import tempfile

    elem_type = payload.get("type", "chart")  # 'chart' or 'table'
    title = payload.get("title", "데이터")
    data = payload.get("data", {})

    prs = Prs()
    prs.slide_width = Inches(10.8)
    prs.slide_height = Inches(7.5)
    blank_layout = prs.slide_layouts[6]  # blank
    slide = prs.slides.add_slide(blank_layout)

    if elem_type == "chart":
        # ── 네이티브 차트 (데이터시트 포함, 더블클릭 편집 가능) ──
        labels = data.get("labels", [])
        values = data.get("values", [])
        chart_type_str = data.get("chartType", "bar")

        cd = ChartData()
        cd.categories = labels
        cd.add_series('평균', [round(v, 2) for v in values])

        ct = XL_CHART_TYPE.COLUMN_CLUSTERED
        if chart_type_str == "horizontalBar":
            ct = XL_CHART_TYPE.BAR_CLUSTERED
        elif chart_type_str == "line":
            ct = XL_CHART_TYPE.LINE

        chart_frame = slide.shapes.add_chart(
            ct, Inches(0.5), Inches(1.2), CHART_SIZE["width"], CHART_SIZE["height"], cd
        )
        chart = chart_frame.chart
        chart.has_legend = False

        try:
            plot = chart.plots[0]
            plot.gap_width = 80
            plot.has_data_labels = True
            dl = plot.data_labels
            dl.number_format = CHART_CONFIG["label_format"]
            dl.font.size = Pt(10)
            dl.font.bold = True
            dl.font.name = FONT_NAME

            colors = data.get("colors", [])
            for pi in range(len(values)):
                try:
                    pt = plot.series[0].points[pi]
                    pt.format.fill.solid()
                    if pi < len(colors) and colors[pi]:
                        hex_c = colors[pi].lstrip('#')
                        pt.format.fill.fore_color.rgb = RGBColor(
                            int(hex_c[0:2],16), int(hex_c[2:4],16), int(hex_c[4:6],16))
                    else:
                        v = values[pi]
                        if v >= 4.5:
                            pt.format.fill.fore_color.rgb = BRAND_COLORS["SUCCESS"]
                        elif v < 3.5:
                            pt.format.fill.fore_color.rgb = BRAND_COLORS["DANGER"]
                        else:
                            pt.format.fill.fore_color.rgb = BRAND_COLORS["BLUE"]
                except: pass
        except: pass

        try:
            plot = chart.plots[0]
            plot.gap_width = CHART_CONFIG["gap_width"]
            chart.value_axis.minimum_scale = 0
            chart.value_axis.maximum_scale = 5
            chart.value_axis.major_unit = 1
            chart.category_axis.tick_labels.font.size = Pt(10)
            chart.category_axis.tick_labels.font.name = FONT_NAME
            chart.value_axis.tick_labels.font.size = Pt(10)
            chart.value_axis.tick_labels.font.name = FONT_NAME
            # 눈금선 제거
            chart.value_axis.has_major_gridlines = True
            chart.value_axis.major_gridlines.format.line.color.rgb = BRAND_COLORS["GRID"]
            chart.value_axis.has_minor_gridlines = False
        except: pass

    elif elem_type == "table":
        # ── 네이티브 표 (복사 붙여넣기 가능) ──
        rows_data = data.get("rows", [])
        headers = data.get("headers", [])
        style = data.get("style", "B")

        n_rows = len(rows_data) + 1  # +헤더
        n_cols = len(headers) if headers else (len(rows_data[0]) if rows_data else 4)
        n_rows = max(n_rows, 2)
        n_cols = max(n_cols, 1)

        tbl_shape = slide.shapes.add_table(
            n_rows, n_cols,
            Inches(0.5), Inches(1.5), TABLE_SIZE["width"], TABLE_SIZE["height"]
        )
        tbl = tbl_shape.table

        # 헤더
        for ci, h in enumerate(headers):
            if ci < n_cols:
                cell = tbl.cell(0, ci)
                cell.text = str(h)
                for p in cell.text_frame.paragraphs:
                    p.alignment = PP_ALIGN.CENTER
                    for r in p.runs:
                        r.font.bold = True
                        r.font.size = Pt(11)
                # 헤더 배경 D9D9D9
                from pptx.oxml.ns import qn
                from lxml import etree
                tc = cell._tc
                tcPr = tc.find(qn('a:tcPr'))
                if tcPr is None:
                    tcPr = etree.SubElement(tc, qn('a:tcPr'))
                solidFill = etree.SubElement(tcPr, qn('a:solidFill'))
                srgb = etree.SubElement(solidFill, qn('a:srgbClr'))
                srgb.set('val', 'F1F5F9') # HEX_GRAY_100 equivalent

        # 데이터 행
        for ri, row in enumerate(rows_data):
            for ci, val in enumerate(row):
                if ci < n_cols and ri + 1 < n_rows:
                    cell = tbl.cell(ri + 1, ci)
                    cell.text = str(val)
                    for p in cell.text_frame.paragraphs:
                        p.alignment = PP_ALIGN.CENTER
                        for r in p.runs:
                            r.font.size = Pt(10)

    # 저장
    out_dir = BASE_DIR / "output"
    out_dir.mkdir(exist_ok=True)
    safe_title = title.replace(' ','_')[:20]
    filename = f"[복사용] {safe_title}.pptx"
    out_path = out_dir / filename
    prs.save(str(out_path))

    safe_fn = quote(filename)
    return FileResponse(
        str(out_path),
        media_type="application/vnd.openxmlformats-officedocument.presentationml.presentation",
        headers={"Content-Disposition": f"attachment; filename*=UTF-8''{safe_fn}"},
    )

@app.post("/api/export_slide")
async def export_slide(payload: dict):
    """차트+표 합쳐서 하나의 슬라이드에 넣은 PPTX 생성 (템플릿 디자인 매칭)"""
    from pptx import Presentation as Prs
    from pptx.util import Inches, Pt, Emu
    from pptx.dml.color import RGBColor
    from pptx.chart.data import ChartData
    from pptx.enum.chart import XL_CHART_TYPE
    from pptx.enum.text import PP_ALIGN
    from urllib.parse import quote

    from modules.ppt_constants import (
        FONT_NAME, BRAND_COLORS,
        CHART_SIZE, TABLE_SIZE,
        CHART_CONFIG
    )

    def _hex_to_rgb(h):
        h = h.lstrip('#')
        return RGBColor(int(h[0:2],16), int(h[2:4],16), int(h[4:6],16))

    def _pastel(h):
        h = h.lstrip('#')
        r, g, b = int(h[0:2],16), int(h[2:4],16), int(h[4:6],16)
        r = int(r * 0.6 + 255 * 0.4)
        g = int(g * 0.6 + 255 * 0.4)
        b = int(b * 0.6 + 255 * 0.4)
        return RGBColor(min(r,255), min(g,255), min(b,255))

    title = payload.get("title", "슬라이드")
    chart_data = payload.get("chart")
    table_data = payload.get("table")

    prs = Prs()
    prs.slide_width = Inches(13.33)
    prs.slide_height = Inches(7.5)
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    # 제목 (세로바 + 텍스트)
    bar_shape = slide.shapes.add_shape(1, Inches(0.4), Inches(0.25), Inches(0.08), Inches(0.35))
    bar_shape.fill.solid()
    bar_shape.fill.fore_color.rgb = BRAND_COLORS["BLUE"]
    bar_shape.line.fill.background()

    txBox = slide.shapes.add_textbox(Inches(0.6), Inches(0.2), Inches(12), Inches(0.45))
    tf = txBox.text_frame
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(16)
    p.font.bold = True
    p.font.name = FONT_NAME
    p.font.color.rgb = RGBColor(0x33, 0x33, 0x33)

    has_chart = chart_data is not None
    has_table = table_data is not None
    if has_chart and has_table:
        chart_top, chart_h = Inches(0.75), Inches(4.0)
        table_top = Inches(5.0)
    elif has_chart:
        chart_top, chart_h = Inches(0.75), Inches(5.8)
        table_top = None
    else:
        chart_top, chart_h = None, None
        table_top = Inches(0.75)

    if has_chart:
        labels = chart_data.get("labels", [])
        values = chart_data.get("values", [])
        colors_hex = chart_data.get("colors", [])
        chart_type_str = chart_data.get("chartType", "bar")
        cd = ChartData()
        cd.categories = labels
        cd.add_series('평균', [round(v, 2) for v in values])
        ct = XL_CHART_TYPE.COLUMN_CLUSTERED
        if chart_type_str == "horizontalBar":
            ct = XL_CHART_TYPE.BAR_CLUSTERED
        elif chart_type_str == "line":
            ct = XL_CHART_TYPE.LINE
        chart_frame = slide.shapes.add_chart(ct, Inches(0.5), chart_top, CHART_SIZE["width"], CHART_SIZE["height"], cd)
        chart_obj = chart_frame.chart
        chart_obj.has_legend = False
        chart_obj.font.name = FONT_NAME
        plot = chart_obj.plots[0]
        plot.has_data_labels = True
        plot.data_labels.font.size = Pt(9)
        plot.data_labels.font.bold = True
        plot.data_labels.font.name = FONT_NAME
        plot.data_labels.number_format = CHART_CONFIG["label_format"]
        if ct != XL_CHART_TYPE.LINE:
            plot.gap_width = CHART_CONFIG["gap_width"]
            series = plot.series[0]
            for i, c_hex in enumerate(colors_hex):
                if i < len(series.points):
                    series.points[i].format.fill.solid()
                    series.points[i].format.fill.fore_color.rgb = _hex_to_rgb(c_hex) if c_hex else BRAND_COLORS["BLUE"]
        va = chart_obj.value_axis
        va.maximum_scale = 5.0
        va.minimum_scale = 0
        va.major_unit = 1.0
        va.has_major_gridlines = True
        va.major_gridlines.format.line.color.rgb = RGBColor(0xDD, 0xDD, 0xDD)
        va.tick_labels.font.size = Pt(8)
        va.tick_labels.font.name = FONT_NAME
        ca = chart_obj.category_axis
        ca.tick_labels.font.size = Pt(8)
        ca.tick_labels.font.name = FONT_NAME

    if has_table:
        headers = table_data.get("headers", [])
        rows_data = table_data.get("rows", [])
        t_colors = chart_data.get("colors", []) if chart_data else []
        n_cols = len(headers)
        n_rows = 1 + len(rows_data)
        row_h = Inches(0.32)
        tbl_shape = slide.shapes.add_table(n_rows, n_cols, Inches(0.5), table_top, TABLE_SIZE["width"], TABLE_SIZE["height"])
        tbl = tbl_shape.table
        for ci, h in enumerate(headers):
            cell = tbl.cell(0, ci)
            cell.text = str(h)
            cell.fill.solid()
            cell.fill.fore_color.rgb = _pastel(t_colors[ci]) if ci < len(t_colors) else RGBColor(0xF0, 0xF0, 0xF0)
            for pp in cell.text_frame.paragraphs:
                pp.alignment = PP_ALIGN.CENTER
                for rr in pp.runs:
                    rr.font.size = Pt(9)
                    rr.font.bold = True
                    rr.font.name = FONT_NAME
                    rr.font.color.rgb = RGBColor(0x33, 0x33, 0x33)
        for ri, row in enumerate(rows_data):
            for ci, val in enumerate(row):
                if ci < n_cols and ri + 1 < n_rows:
                    cell = tbl.cell(ri + 1, ci)
                    cell.text = str(val)
                    cell.fill.solid()
                    cell.fill.fore_color.rgb = RGBColor(0xFF, 0xFF, 0xFF)
                    for pp in cell.text_frame.paragraphs:
                        pp.alignment = PP_ALIGN.CENTER
                        for rr in pp.runs:
                            rr.font.size = Pt(9)
                            rr.font.name = FONT_NAME

    out_dir = BASE_DIR / "output"
    out_dir.mkdir(exist_ok=True)
    safe_title = title.replace(' ','_')[:20]
    filename = f"{safe_title}.pptx"
    out_path = out_dir / filename
    prs.save(str(out_path))
    safe_fn = quote(filename)
    return FileResponse(
        str(out_path),
        media_type="application/vnd.openxmlformats-officedocument.presentationml.presentation",
        headers={"Content-Disposition": f"attachment; filename*=UTF-8''{safe_fn}"},
    )

@app.post("/api/upload_template/{session_id}")
async def upload_template(session_id: str, file: UploadFile = File(...)):
    """PPT 템플릿 업로드 — 세션에 저장"""
    if session_id not in sessions:
        raise HTTPException(404, "세션 없음")
    if not file.filename.lower().endswith('.pptx'):
        raise HTTPException(400, ".pptx 파일만 허용됩니다")

    tpl_dir = BASE_DIR / "templates"
    tpl_dir.mkdir(exist_ok=True)
    tpl_path = tpl_dir / f"custom_{session_id}.pptx"

    content = await file.read()
    with open(tpl_path, 'wb') as f:
        f.write(content)

    sessions[session_id]["template_path"] = str(tpl_path)

    # 템플릿 슬라이드 수 확인
    try:
        from pptx import Presentation as Prs
        prs = Prs(str(tpl_path))
        slide_count = len(prs.slides)
    except Exception:
        slide_count = 0

    return {"ok": True, "filename": file.filename, "slide_count": slide_count}


@app.post("/api/modify/{session_id}")

async def modify(session_id: str, request: dict):
    """수정 요청 → AI 해석 → 재생성"""
    if session_id not in sessions:
        raise HTTPException(404, "세션 없음")
    
    session = sessions[session_id]
    user_request = request.get("message", "")
    
    modification = await ai_engine.interpret_modification(
        user_request, session["summary"], session.get("preview", [])
    )
    
    data = session["data"]
    matched = session.get("matched_sample")
    multi_result = session.get("multi_result", {})
    output_path = generate_report(data, config, sample_pattern=matched, multi_result=multi_result)
    session["output_path"] = output_path
    
    preview = generate_preview(output_path)
    session["preview"] = preview
    
    return {
        "success": True,
        "modification": modification,
        "slide_count": len(preview),
        "preview": preview,
    }


@app.get("/api/download/{session_id}")
async def download(session_id: str):
    """PPT 다운로드"""
    if session_id not in sessions or "output_path" not in sessions[session_id]:
        raise HTTPException(404, "파일 없음")
    
    path = sessions[session_id]["output_path"]
    return FileResponse(
        path,
        media_type="application/vnd.openxmlformats-officedocument.presentationml.presentation",
        filename=os.path.basename(path),
    )


@app.post("/api/samples/add")
async def add_sample(file: UploadFile = File(...)):
    """새 샘플 추가 → 자동 패턴 분석"""
    sample_path = BASE_DIR / "samples" / file.filename
    with open(sample_path, "wb") as f:
        shutil.copyfileobj(file.file, f)
    
    pattern = sample_analyzer.analyze_single(str(sample_path))
    sample_analyzer.save_patterns()
    
    return {"success": True, "pattern": pattern}


@app.get("/api/samples")
async def list_samples():
    """샘플 패턴 목록"""
    return sample_analyzer.get_patterns()


@app.get("/api/sheets")
async def get_sheets_api(path: str):
    """업로드된 엑셀의 시트 목록 (미리보기용)"""
    try:
        sheets = get_available_sheets(path)
        return {"sheets": sheets}
    except:
        return {"sheets": []}


@app.get("/api/status")
async def api_status():
    """AI 연결 상태"""
    return {
        "ai_enabled": ai_engine.enabled,
        "api_key_set": bool(config.get("gemini", {}).get("api_key")),
        "model": config.get("gemini", {}).get("model", "gemini-2.0-flash"),
        "sample_count": len(sample_analyzer.get_patterns()),
    }


if __name__ == "__main__":
    import uvicorn
    uvicorn.run(app, host="0.0.0.0", port=8080)
