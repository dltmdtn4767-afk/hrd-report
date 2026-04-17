"""Template PPTX full analysis"""
import sys, io
sys.stdout = io.TextIOWrapper(sys.stdout.buffer, encoding='utf-8', errors='replace')

from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE
import json

prs = Presentation("templates/template.pptx")
print(f"Size: {prs.slide_width/914400:.1f}x{prs.slide_height/914400:.1f} in, Slides: {len(prs.slides)}")

result = []
for si, slide in enumerate(prs.slides):
    info = {"index": si+1, "texts": [], "tables": [], "charts": [], "images": [], "shapes_detail": []}
    print(f"\n--- Slide {si+1} ---")
    for sh in slide.shapes:
        sd = {"name": sh.name, "type": str(sh.shape_type),
              "pos": f"{sh.left/914400:.1f},{sh.top/914400:.1f}",
              "size": f"{sh.width/914400:.1f}x{sh.height/914400:.1f}"}
        
        if sh.has_text_frame:
            for p in sh.text_frame.paragraphs:
                t = p.text.strip()
                if not t: continue
                fi = {}
                if p.runs:
                    r = p.runs[0]
                    fi = {"name": str(r.font.name or ""), "size": str(r.font.size or ""),
                          "bold": r.font.bold}
                    try: fi["color"] = str(r.font.color.rgb)
                    except: fi["color"] = ""
                print(f"  TXT [{sh.name}]: {t[:60]}  font={fi}")
                info["texts"].append({"shape": sh.name, "text": t, "font": fi})
        
        if sh.has_table:
            tbl = sh.table
            rows = []
            print(f"  TABLE [{sh.name}]: {len(tbl.rows)}r x {len(tbl.columns)}c")
            for ri, row in enumerate(tbl.rows):
                rd = []
                for ci, cell in enumerate(row.cells):
                    ct = cell.text.strip()
                    fc = ""
                    try: fc = str(cell.fill.fore_color.rgb)
                    except: pass
                    cf = {}
                    try:
                        r = cell.text_frame.paragraphs[0].runs[0]
                        cf = {"name": str(r.font.name or ""), "size": str(r.font.size or ""),
                              "bold": r.font.bold}
                        try: cf["color"] = str(r.font.color.rgb)
                        except: cf["color"] = ""
                    except: pass
                    rd.append({"text": ct, "bg": fc, "font": cf})
                    if ri < 4:
                        print(f"    [{ri},{ci}] '{ct[:25]}' bg={fc} f={cf}")
                rows.append(rd)
            info["tables"].append({"rows": len(tbl.rows), "cols": len(tbl.columns), "data": rows})
        
        if sh.has_chart:
            ct = str(sh.chart.chart_type)
            print(f"  CHART [{sh.name}]: {ct}")
            info["charts"].append({"type": ct, "name": sh.name})
        
        if sh.shape_type == MSO_SHAPE_TYPE.PICTURE:
            print(f"  IMG [{sh.name}]: {sh.image.content_type}")
            info["images"].append({"name": sh.name, "type": sh.image.content_type})
        
        info["shapes_detail"].append(sd)
    result.append(info)

with open("template_analysis.json", "w", encoding="utf-8") as f:
    json.dump(result, f, ensure_ascii=False, indent=2, default=str)
print("\n=== DONE: template_analysis.json ===")
